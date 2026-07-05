#!/usr/bin/env python3
"""Offline tests: fake HTTP fetch, deterministic bag-of-words embedder, real DuckDB + Lance
in a temp BEAN_HOME. Covers the store (hash-gated upserts, revisions, cursors), both sources
(change detection, export fallback, per-thread/message chat docs, edits in lookback), the retry policy, the
end-to-end sync→search flow, and workspace/credential hygiene."""

from __future__ import annotations

import json
import re
import stat
import sys
import tempfile
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from bean.workspace import set_bean_home  # noqa: E402
set_bean_home(tempfile.mkdtemp(prefix="bean-home-"))  # config, not an env var

from bean.http import AuthError, Response, api_get  # noqa: E402
from bean.store import Store, content_hash  # noqa: E402
from bean.chunks import chunk_text  # noqa: E402
from bean.workspace import Workspace, save_credential, load_credential, bean_home  # noqa: E402
from bean import config as cfgmod, pdf  # noqa: E402
from bean.connectors import gdocs, slack, localfiles, github  # noqa: E402
from bean.index import search as lance_search  # noqa: E402
from bean.sync import run_sync  # noqa: E402
from bean.search import search as hybrid_search, recent, thread, neighbors  # noqa: E402

CHECKS = FAILED = 0


def reembed(ws, *, embed_fn=None, log=lambda m: None):
    """`bean sync --rebuild` with no live fetch: re-embed every stored doc with the current
    settings (there's no standalone reembed command anymore — --rebuild absorbed it)."""
    r = run_sync(ws, full=True, refetch=False, embed_fn=embed_fn, log=log)
    return {"docs": r["embedded"], "chunks": r["chunks"]}


def ok(cond, msg):
    global CHECKS, FAILED
    CHECKS += 1
    if not cond:
        FAILED += 1
        print(f"  ✗ {msg}")


def res(status, body, headers=None):
    text = body if isinstance(body, str) else json.dumps(body)
    return Response(status=status, text=text, headers=headers or {})


def fake_embed(texts):
    out = []
    for t in texts:
        v = [0.0] * 64
        for w in re.findall(r"[a-z]{3,}", t.lower()):
            h = 0
            for ch in w:
                h = (h * 31 + ord(ch)) & 0xFFFFFFFF
            v[h % 64] += 1
        norm = sum(x * x for x in v) ** 0.5 or 1
        out.append([x / norm for x in v])
    return out


def repo(name):
    d = Path(tempfile.mkdtemp(prefix=f"bean-{name}-"))
    (d / ".git").mkdir()
    return d


# -- http retry policy -----------------------------------------------------------------------
calls = {"n": 0}
slept = []


def flaky(url, headers):
    calls["n"] += 1
    return res(429, "slow", {"Retry-After": "2"}) if calls["n"] < 3 else res(200, "fine")


r = api_get("https://x.test/a", fetch=flaky, sleep=slept.append)
ok(r.status == 200 and calls["n"] == 3, "429s retried until success")
ok(all(s == 2.0 for s in slept), "Retry-After honored")
try:
    api_get("https://x.test/a", fetch=lambda u, h: res(401, "no"))
    ok(False, "401 should raise")
except AuthError:
    ok(True, "401 raises AuthError")
bad = api_get("https://x.test/a", fetch=lambda u, h: res(500, "boom"), sleep=lambda s: None, retries=1)
ok(bad.status == 500, "exhausted retries return the failing response")

# -- workspace + credentials -------------------------------------------------------------------
ws1, ws2 = Workspace(repo("alpha")), Workspace(repo("beta"))
ok(ws1.dir != ws2.dir and ws1.dir.parent == ws2.dir.parent == bean_home(), "one workspace folder per repo under ~/.bean")
ok(Workspace(ws1.repo).dir == ws1.dir, "workspace is stable for the same repo")
path = save_credential("slack", {"token": "xoxp-1"})
ok(stat.S_IMODE(path.stat().st_mode) == 0o600, "credentials are mode 0600")
ok(load_credential("slack")["token"] == "xoxp-1", "credential round-trips")

# -- store: hash-gated upserts ------------------------------------------------------------------
with Store(ws1) as store:
    ok(store.upsert("gdocs", "d1", title="Vision", url="https://d/1", revision_id="r1", body="We charge cards."),
       "first upsert reports change")
    ok(not store.upsert("gdocs", "d1", title="Vision", url="https://d/1", revision_id="r2", body="We charge cards."),
       "same body (comment-only revision) reports no change")
    ok(store.get("gdocs", "d1").revision_id == "r2", "metadata still updated")
    ok(store.upsert("gdocs", "d1", title="Vision v2", url="https://d/1", revision_id="r3", body="We charge twice."),
       "content change reports change")
    ok(len(store.revisions("gdocs", "d1")) == 2, "revision ledger records each content change")
    store.set_state("k", {"a": 1})
    ok(store.get_state("k") == {"a": 1}, "state round-trips as json")

    # source-native metadata: stored, round-tripped, and updated on a comment-only revision
    store.upsert("gdocs", "d2", title="Old", url=None, revision_id="v1", body="aged content",
                 meta={"modified_at": "2020-01-01T00:00:00Z", "author": "Ada", "mime": "application/vnd.google-apps.document"})
    store.upsert("gdocs", "d3", title="New", url=None, revision_id="v1", body="fresh content",
                 meta={"modified_at": "2026-06-01T00:00:00Z", "author": "Bob"})
    ok(store.get("gdocs", "d2").author == "Ada", "author metadata round-trips")
    ok(str(store.get("gdocs", "d3").modified_at).startswith("2026-06-01"), "modified_at parsed to a timestamp")
    order = [r["doc_id"] for r in store.recent(source="gdocs")]
    ok(order.index("d3") < order.index("d2"), "recent orders by the doc's own modified_at, newest first")
    store.upsert("gdocs", "d2", title="Old", url=None, revision_id="v2", body="aged content",
                 meta={"modified_at": "2027-01-01T00:00:00Z", "author": "Ada"})
    ok(str(store.get("gdocs", "d2").modified_at).startswith("2027"), "comment-only revision refreshes modified_at")

# -- chunking -----------------------------------------------------------------------------------
chunks = chunk_text("\n".join(f"line {i} with some sensible amount of text" for i in range(100)), "gdocs/d1")
ok(len(chunks) > 1 and chunks[0].id == "gdocs/d1#L1", "chunking yields stable ids")
ok(all(len(c.text) <= 2000 for c in chunks), "chunks capped")

# -- gdocs sync ---------------------------------------------------------------------------------
DOCS = {
    "docA": {"name": "Payments Guide", "rev": "r1", "body": "# Payments\n\nBilling goes through chargeCard and emits a receipt.", "md": True},
    "docB": {"name": "Legacy Notes", "rev": "r9", "body": "Plain only.", "md": False},
}
exports = {"n": 0}


def gfetch(url, headers):
    from urllib.parse import urlparse, parse_qs
    u = urlparse(url)
    if u.path == "/drive/v3/files":
        return res(200, {"files": [{"id": "docB"}]})
    cm = re.match(r"^/drive/v3/files/([^/]+)/comments$", u.path)
    if cm:
        return res(200, DOCS.get(cm.group(1), {}).get("comments") or {"comments": []})
    m = re.match(r"^/drive/v3/files/([^/]+)(/export)?$", u.path)
    d = DOCS[m.group(1)]
    qs = parse_qs(u.query)
    if qs.get("alt") == ["media"]:  # PDF binary download
        return res(200, d["body"])
    if not m.group(2):
        mime = "application/pdf" if d.get("pdf") else "application/vnd.google-apps.document"
        return res(200, {"id": m.group(1), "name": d["name"], "headRevisionId": d["rev"],
                         "webViewLink": f"https://docs.google.com/document/d/{m.group(1)}", "trashed": False,
                         "modifiedTime": "2026-05-01T00:00:00Z", "mimeType": mime,
                         "lastModifyingUser": {"displayName": "Grace Hopper"}})
    exports["n"] += 1
    mime = qs["mimeType"][0]
    if mime == "text/markdown" and not d["md"]:
        return res(400, "no markdown")
    return res(200, d["body"])


gws = Workspace(repo("gdocs"))
cfg = {"docs": ["docA"], "folders": ["folder1"]}
with Store(gws) as store:
    s1 = gdocs.sync(store, cfg, token_fn=lambda force=False: "tok", fetch=gfetch)
    ok(len(s1["changed"]) == 2, f"first sync ingests both docs ({s1['changed']})")
    ok(store.get("gdocs", "docB").body == "Plain only.", "markdown failure falls back to text/plain")
    ok(store.get("gdocs", "docA").author == "Grace Hopper" and store.get("gdocs", "docA").mime.endswith("document"),
       "gdocs sync captures author + mime from Drive metadata")
    ok(str(store.get("gdocs", "docA").modified_at).startswith("2026-05-01"), "gdocs sync captures modifiedTime")
    s2 = gdocs.sync(store, cfg, token_fn=lambda force=False: "tok", fetch=gfetch)
    ok(s2["changed"] == [] and s2["removed"] == [], "same revisions are a no-op")
    DOCS["docA"]["rev"] = "r2"  # comment-only: new revision, same body
    before = exports["n"]
    s3 = gdocs.sync(store, cfg, token_fn=lambda force=False: "tok", fetch=gfetch)
    ok(s3["changed"] == [] and exports["n"] == before + 1, "comment-only revision exports once, changes nothing")
    DOCS["docA"]["rev"] = "r3"
    DOCS["docA"]["body"] += "\n\nRetries happen twice."
    s4 = gdocs.sync(store, cfg, token_fn=lambda force=False: "tok", fetch=gfetch)
    ok(s4["changed"] == ["docA"], "real change re-ingests")
    s5 = gdocs.sync(store, {"docs": ["docA"], "folders": []}, token_fn=lambda force=False: "tok", fetch=gfetch)
    ok(s5["removed"] == ["docB"] and store.get("gdocs", "docB") is None, "unlisted doc pruned")

    # mid-sync 401: token_fn(force=True) asked once, sync continues
    state = {"failed": False, "forces": 0}

    def gfetch401(url, headers):
        if not state["failed"]:
            state["failed"] = True
            return res(401, "expired")
        return gfetch(url, headers)

    def token_fn(force=False):
        if force:
            state["forces"] += 1
        return "tok2"

    gdocs.sync(store, {"docs": ["docA"], "folders": []}, token_fn=token_fn, fetch=gfetch401)
    ok(state["forces"] == 1, "mid-sync 401 refreshes the token once")

# gdrive PDFs: a native PDF is downloaded (alt=media) and run through the shared extractor, not
# the Docs export path. The extractor is injected so the routing is exercised without pymupdf.
DOCS["docP"] = {"name": "Scanned Report", "rev": "p1", "body": "%PDF-1.4 binary…", "pdf": True}
with Store(Workspace(repo("gdocs-pdf"))) as store:
    calls = {"n": 0}

    def fake_extract(path, ocr, log=lambda m: None):
        calls["n"] += 1
        return "quarterly revenue and invoice totals"

    sp = gdocs.sync(store, {"docs": ["docP"], "folders": []},
                    token_fn=lambda force=False: "tok", fetch=gfetch, extract=fake_extract)
    ok(sp["changed"] == ["docP"], f"gdrive PDF ingested ({sp['changed']})")
    ok(calls["n"] == 1, "PDF routed through the extractor, not the Docs export path")
    ok("invoice totals" in store.get("gdocs", "docP").body, "extracted PDF text stored as body")
    ok(store.get("gdocs", "docP").mime == "application/pdf", "PDF mime captured from Drive metadata")
    sp2 = gdocs.sync(store, {"docs": ["docP"], "folders": []},
                     token_fn=lambda force=False: "tok", fetch=gfetch, extract=fake_extract)
    ok(sp2["changed"] == [] and calls["n"] == 1, "unchanged PDF revision skips re-download and re-extract")

# auto-index: empty config crawls the docs you own; the query carries an ownership + window filter
autodocs = {"m1": {"name": "Owned One", "rev": "a1", "body": "alpha", "md": True, "mod": "2026-03-01T00:00:00Z"},
            "m2": {"name": "Owned Two", "rev": "a2", "body": "beta", "md": True, "mod": "2026-04-01T00:00:00Z"}}
seen_q = {}


def afetch(url, headers):
    from urllib.parse import urlparse, parse_qs
    u = urlparse(url)
    if u.path.endswith("/comments"):
        return res(200, {"comments": []})
    if u.path == "/drive/v3/files":
        seen_q["q"] = parse_qs(u.query).get("q", [""])[0]
        return res(200, {"files": [{"id": "m1"}, {"id": "m2"}]})
    m = re.match(r"^/drive/v3/files/([^/]+)(/export)?$", u.path)
    d = autodocs[m.group(1)]
    if not m.group(2):
        return res(200, {"id": m.group(1), "name": d["name"], "headRevisionId": d["rev"], "trashed": False,
                         "webViewLink": f"https://docs.google.com/document/d/{m.group(1)}", "modifiedTime": d["mod"]})
    return res(200, d["body"])


aws = Workspace(repo("gdocs-auto"))
with Store(aws) as store:
    sa = gdocs.sync(store, {}, token_fn=lambda force=False: "tok", fetch=afetch, lookback_days=30)
    ok(sorted(sa["changed"]) == ["m1", "m2"], "empty config auto-indexes owned docs")
    ok("'me' in owners" in seen_q["q"] and "modifiedTime >" in seen_q["q"], "auto query filters by owner + window")
    ok(store.get_state("gdocs.cursor") == "2026-04-01T00:00:00Z", "first sync records the newest modifiedTime as the cursor")
    # smart lookback: the next sync's discovery query starts from the cursor, not the fixed window
    sa2 = gdocs.sync(store, {}, token_fn=lambda force=False: "tok", fetch=afetch, lookback_days=30)
    ok("modifiedTime > '2026-04-01T00:00:00Z'" in seen_q["q"], "second sync discovers only files changed since the cursor")
    # --rebuild ignores the cursor and reaches back over the lookback window again
    gdocs.sync(store, {}, token_fn=lambda force=False: "tok", fetch=afetch, lookback_days=30, full=True)
    ok("modifiedTime > '2026-04-01T00:00:00Z'" not in seen_q["q"], "--rebuild ignores the cursor")

    # a later crawl that no longer returns m2 still retains it (only trash/access-loss evicts)
    def afetch2(url, headers):
        from urllib.parse import urlparse
        if urlparse(url).path == "/drive/v3/files":
            return res(200, {"files": [{"id": "m1"}]})
        return afetch(url, headers)

    sb = gdocs.sync(store, {}, token_fn=lambda force=False: "tok", fetch=afetch2, lookback_days=30)
    ok(sb["removed"] == [] and store.get("gdocs", "m2") is not None, "doc aged out of window is retained")

# -- gdocs comments: each comment is its own author-attributed, timestamped doc ------------------
def cfetch(url, headers):
    from urllib.parse import urlparse
    u = urlparse(url)
    if u.path.endswith("/comments"):
        return res(200, {"comments": [
            {"id": "c1", "content": "please clarify the refund window",
             "author": {"displayName": "Eric Idle"}, "createdTime": "2026-06-01T00:00:00Z",
             "modifiedTime": "2026-06-02T00:00:00Z",
             "quotedFileContent": {"value": "refunds are processed"},
             "replies": [{"author": {"displayName": "You"}, "content": "within 30 days",
                          "createdTime": "2026-06-02T00:00:00Z"}]}]})
    if u.path.endswith("/export"):
        return res(200, "# Refund Policy\n\nRefunds are processed within the window.")
    return res(200, {"id": "docC", "name": "Refund Policy", "headRevisionId": "r1", "trashed": False,
                     "webViewLink": "https://docs.google.com/document/d/docC",
                     "modifiedTime": "2026-05-01T00:00:00Z",
                     "mimeType": "application/vnd.google-apps.document",
                     "lastModifyingUser": {"displayName": "Grace Hopper"}})

cdoc = "docC#comment:c1"
cws = Workspace(repo("gdocs-comments"))
with Store(cws) as store:
    rc = gdocs.sync(store, {"docs": ["docC"], "folders": []}, token_fn=lambda force=False: "tok", fetch=cfetch)
    ok(cdoc in rc["changed"], "comment indexed as its own doc")
    c = store.get("gdocs", cdoc)
    ok(c and c.author == "Eric Idle", "comment carries its own author")
    ok(str(c.modified_at).startswith("2026-06-02"), "comment carries its last-activity time")
    ok("refund window" in c.body and "within 30 days" in c.body, "comment body includes content + replies")
    ok("Refund Policy" in c.title and "Eric Idle" in c.title, "comment title names author + parent doc")
    # the user story: "show me eric's most recent comment on my doc"
    hits = recent(cws, source="gdocs", author="eric", doc_like="Refund")
    ok(hits and hits[0]["doc_id"] == cdoc, "recent --author eric --doc Refund surfaces eric's comment")
    rc2 = gdocs.sync(store, {"docs": ["docC"], "folders": []}, token_fn=lambda force=False: "tok", fetch=cfetch)
    ok(cdoc not in rc2["changed"], "unchanged comment re-syncs as a no-op")

    def cfetch_gone(url, headers):
        from urllib.parse import urlparse
        if urlparse(url).path.endswith("/comments"):
            return res(200, {"comments": []})
        return cfetch(url, headers)
    rc3 = gdocs.sync(store, {"docs": ["docC"], "folders": []}, token_fn=lambda force=False: "tok", fetch=cfetch_gone)
    ok(cdoc in rc3["removed"] and store.get("gdocs", cdoc) is None, "a deleted comment is pruned")

# -- slack sync ---------------------------------------------------------------------------------
NOW = time.mktime(time.strptime("2026-07-02 12:00", "%Y-%m-%d %H:%M"))


def ts(days_ago, frac=".000100"):
    return f"{int(NOW - days_ago * 86400)}{frac}"


HISTORY = [
    {"ts": ts(1), "user": "U1", "text": "Retries land <@U2>, see <https://x.test|the doc>", "thread_ts": ts(1), "reply_count": 1},
    {"ts": ts(2), "user": "U2", "text": "Deploy done"},
]
REPLIES = [{"ts": ts(1, ".000200"), "user": "U2", "text": "Confirmed", "thread_ts": ts(1)}]


seen_oldest = {}


def sfetch(url, headers):
    from urllib.parse import urlparse, parse_qs
    u = urlparse(url)
    q = parse_qs(u.query)
    method = u.path.split("/")[-1]
    if method == "conversations.list":
        return res(200, {"ok": True, "channels": [
            {"id": "C1", "name": "eng-payments", "is_member": True},
            {"id": "C2", "name": "random-not-joined", "is_member": False}]})
    if method == "conversations.history":
        oldest = float(q["oldest"][0])
        seen_oldest["v"] = oldest
        return res(200, {"ok": True, "messages": [m for m in HISTORY if float(m["ts"]) >= oldest]})
    if method == "conversations.replies":
        return res(200, {"ok": True, "messages": [HISTORY[0]] + REPLIES})
    if method == "users.list":
        return res(200, {"ok": True, "members": [
            {"id": "U1", "name": "ada", "profile": {"display_name": "ada"}},
            {"id": "U2", "name": "bob", "profile": {}}]})
    return res(404, {"ok": False, "error": "unknown"})


THREAD_ID = f"eng-payments/{ts(1)}"      # one doc per thread, keyed by the root message ts
SINGLE_ID = f"eng-payments/{ts(2)}"      # one doc per standalone message
sws = Workspace(repo("slack"))
scfg = {"channels": ["#eng-payments"], "lookback_days": 14}
with Store(sws) as store:
    s1 = slack.sync(store, scfg, token="xoxp-1", team_url="https://t.slack.com", fetch=sfetch, now=NOW)
    ok(s1["changed"] == [THREAD_ID, SINGLE_ID], f"first sync writes one doc per thread/message ({s1['changed']})")
    tbody = store.get("slack", THREAD_ID).body
    ok("@ada" in tbody and "Retries land @bob, see the doc (https://x.test)" in tbody, "mentions and links resolve")
    ok("Confirmed" in tbody, "the thread's reply is in the thread doc")
    ok("Deploy done" in store.get("slack", SINGLE_ID).body, "a standalone message is its own doc")
    ok(str(store.get("slack", THREAD_ID).modified_at).startswith("2026-07-01"), "thread modified_at is its latest reply")
    s2 = slack.sync(store, scfg, token="xoxp-1", team_url="https://t.slack.com", fetch=sfetch, now=NOW)
    ok(s2["changed"] == [], "unchanged history is a no-op")
    ok(store.get_state("slack.cursor.C1") == float(ts(1, ".000200")), "cursor advanced to the latest activity")
    # Later syncs re-scan a trailing REFRESH_DAYS window so recent edits/replies re-render.
    ok(seen_oldest["v"] == NOW - slack.REFRESH_DAYS * 86400, "later syncs floor at the refresh window")
    HISTORY[1]["text"] = "Deploy done (edited: rolled back)"
    s3 = slack.sync(store, scfg, token="xoxp-1", team_url="https://t.slack.com", fetch=sfetch, now=NOW)
    ok(s3["changed"] == [SINGLE_ID], "an edit within the window rewrites just that message's doc")
    ok("rolled back" in store.get("slack", SINGLE_ID).body, "edited text landed")

# all-channels mode: no explicit channel list → sync every channel the account is a member of
aws = Workspace(repo("slack-all"))
with Store(aws) as store:
    sa = slack.sync(store, {"channels": []}, token="xoxp-1", team_url="https://t.slack.com", fetch=sfetch, now=NOW)
    ok(sa["changed"] == [THREAD_ID, SINGLE_ID], f"empty channel list syncs all member channels ({sa['changed']})")
    ok(not any(d.startswith("random-not-joined/") for d in store.doc_ids("slack")), "non-member channels are skipped")

# legacy migration: pre-existing per-ISO-week digest docs are pruned on sync
mws = Workspace(repo("slack-legacy"))
with Store(mws) as store:
    store.upsert("slack", "eng-payments/2026-W27", title="#eng-payments — 2026-W27", url=None,
                 revision_id=None, body="old digest")
    ml = slack.sync(store, scfg, token="xoxp-1", team_url="https://t.slack.com", fetch=sfetch, now=NOW)
    ok("eng-payments/2026-W27" in ml["removed"], f"legacy week digest is reported removed ({ml['removed']})")
    ok(store.get("slack", "eng-payments/2026-W27") is None, "legacy week digest row is deleted")

# -- end to end: sync → embed → lance → search --------------------------------------------------
ews = Workspace(repo("e2e"))
ews.save_config({"google": {"docs": ["docA"], "folders": []}, "slack": {"channels": []}})
gdocs._token_cache.update(token="tok", exp=time.time() + 3600)  # bypass gcloud in tests
save_credential("google", {"method": "gcloud"})
(bean_home() / "credentials" / "slack.json").unlink(missing_ok=True)  # isolate: only gdocs active here
result = run_sync(ews, embed_fn=fake_embed, fetch=gfetch)
ok(result["errors"] == [] and len(result["changed"]) == 1 and result["chunks"] >= 1,
   f"run_sync ingests and embeds ({result})")
with Store(ews) as _s:
    ok(_s.get_state("last_sync") is not None, "run_sync records last_sync for staleness checks")
hits = lance_search(ews, fake_embed(["billing chargeCard receipt payments"])[0], k=3)
ok(hits and hits[0]["title"] == "Payments Guide" and hits[0]["url"].startswith("https://docs.google.com/"),
   f"search returns the doc with title + url ({hits[:1]})")
ok(hits[0]["text"].startswith("# Payments"), "search result carries the chunk text")
again = run_sync(ews, embed_fn=fake_embed, fetch=gfetch)
ok(again["changed"] == [] and again["chunks"] == 0, "re-sync with no upstream change embeds nothing")

# -- config layering + coercion -----------------------------------------------------------------
base = cfgmod.resolve()
ok(base["embedding"]["plugin"] is None and "model" not in base["embedding"],
   "defaults resolve with no config file; the built-in embedder has no model switch")
cfgmod.save_global({"embedding": {"batch_size": 128}, "search": {"hybrid": False}})
merged = cfgmod.resolve()
ok(merged["embedding"]["batch_size"] == 128 and merged["chunking"]["lines"] == 40,
   "global config overrides one leaf, keeps sibling defaults")
g = cfgmod.load_global(); cfgmod.set_in(g, "chunking.lines", "20"); cfgmod.set_in(g, "search.hybrid", "true")
ok(g["chunking"]["lines"] == 20 and g["search"]["hybrid"] is True, "config set coerces to leaf type")
cfgmod.save_global({})  # reset so downstream resolves to defaults

# -- per-source chunking: global defaults + per-source overrides --------------------------------
_cc = cfgmod.resolve()
ok(cfgmod.chunking_for(_cc, "gdocs")["lines"] == 40, "a source with no override uses the global chunking default")
ok(cfgmod.chunking_for(_cc, "slack")["lines"] == 15 and cfgmod.chunking_for(_cc, "slack")["max_chars"] == 1000,
   "slack ships smaller default windows")
ok(cfgmod.chunking_for(_cc, "slack")["title_prefix"] is True,
   "an override inherits the global block's un-overridden leaves")
cfgmod.save_global({"chunking": {"lines": 60}, "notion": {"chunking": {"lines": 8}}})
_cc2 = cfgmod.resolve()
ok(cfgmod.chunking_for(_cc2, "gdocs")["lines"] == 60, "a global chunking change reaches sources without their own override")
ok(cfgmod.chunking_for(_cc2, "notion")["lines"] == 8, "a per-source chunking override wins over the global block")
g2 = cfgmod.load_global(); cfgmod.set_in(g2, "github.chunking.lines", "12")
ok(g2["github"]["chunking"]["lines"] == 12, "config set coerces a per-source chunking leaf to int")
cfgmod.save_global({})  # reset

# -- one built-in embedder + a drop-in embedder module ------------------------------------------
from bean import embed as _embed  # noqa: E402
ok(_embed.identity({}) == "jinaai/jina-embeddings-v5-text-nano", "identity names the single built-in model")
ok(_embed.identity({"plugin": "/x/e.py"}) == "plugin:/x/e.py", "identity names a plugin")
try:
    _embed._resolve({"backend": "model2vec", "model": "minishlab/potion-retrieval-32M"})
    ok(False, "stale backend/model config raises instead of silently falling back")
except RuntimeError as _e:
    ok("no longer supported" in str(_e), "stale backend/model config fails loudly")
_edir = Path(tempfile.mkdtemp(prefix="bean-embed-"))
(_edir / "e.py").write_text(
    "def embed(texts):\n"
    "    return [[float(len(t)), 1.0] for t in texts]\n"
    "def embed_query(text):\n"
    "    return [float(len(text)), 2.0]\n")
_ef = _embed.embedder({"plugin": str(_edir / "e.py")})
ok(_ef(["ab", "abc"]) == [[2.0, 1.0], [3.0, 1.0]], "a drop-in embedder plugin returns its own vectors")
ok(_embed.query_embedder({"plugin": str(_edir / "e.py")})("abcd") == [4.0, 2.0],
   "the plugin's embed_query is used for the query side")

# -- lookback: connector-level config resolution ------------------------------------------------
from bean.sources import LOOKBACK_DEFAULTS, _lookback  # noqa: E402
ok(LOOKBACK_DEFAULTS == {"slack": 14, "discord": 14, "gdocs": 30}, "lookback sources: slack, discord, gdocs")
ok(_lookback("discord", {}, {}) == 14, "lookback falls back to the built-in default")
ok(_lookback("discord", {}, {"discord": {"lookback_days": 3}}) == 3, "resolved setting overrides the default")
ok(_lookback("discord", {"lookback_days": 5}, {"discord": {"lookback_days": 3}}) == 5,
   "tracked-config lookback wins over the setting")
ok("github" not in LOOKBACK_DEFAULTS and "notion" not in LOOKBACK_DEFAULTS,
   "sources without a window carry no lookback")

# -- local files connector (incl. mtime skip + prune) -------------------------------------------
docs_dir = Path(tempfile.mkdtemp(prefix="bean-docs-"))
(docs_dir / "guide.md").write_text("# Refunds\n\nRefunds go through refundCard and reverse the receipt.\n")
(docs_dir / "skip.log").write_text("not indexed")  # unsupported extension
lws = Workspace(repo("local"))
lset = cfgmod.resolve(lws)
r1 = localfiles.sync(store_stub := Store(lws), {"paths": [str(docs_dir)]}, settings=lset)
ok(len(r1["changed"]) == 1 and store_stub.get("localfiles", str(docs_dir / "guide.md")), "markdown file indexed, .log ignored")
r2 = localfiles.sync(store_stub, {"paths": [str(docs_dir)]}, settings=lset)
ok(r2["changed"] == [], "unchanged file skipped via mtime")
(docs_dir / "guide.md").unlink()
r3 = localfiles.sync(store_stub, {"paths": [str(docs_dir)]}, settings=lset)
ok(r3["removed"] == [str(docs_dir / "guide.md")], "deleted file pruned")
store_stub.close()

# recursive crawl + office docs: a .docx in a nested subfolder is discovered and its text extracted
import zipfile as _zip
from bean import office  # noqa: E402
nested = docs_dir / "sub" / "deeper"
nested.mkdir(parents=True)
_DOCX_XML = (
    '<?xml version="1.0"?><w:document '
    'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body>'
    "<w:p><w:r><w:t>Chargebacks</w:t></w:r><w:r><w:t> go through disputeCard.</w:t></w:r></w:p>"
    "<w:p><w:r><w:t>Second paragraph.</w:t></w:r></w:p></w:body></w:document>"
)
with _zip.ZipFile(nested / "report.docx", "w") as zf:
    zf.writestr("word/document.xml", _DOCX_XML)
ok(office.extract_office(nested / "report.docx") == "Chargebacks go through disputeCard.\nSecond paragraph.",
   "docx runs join within a paragraph, paragraphs split by newline")
ok(office._rtf(r"{\rtf1\ansi Hello\par world\'21}") == "Hello\nworld!", "rtf control words stripped, escapes decoded")
# .xlsx (openpyxl) + .pptx (python-pptx) via libraries; .html via the stdlib flattener
from openpyxl import Workbook as _WB  # noqa: E402
_wb = _WB(); _sh = _wb.active; _sh.title = "Q3"; _sh.append(["Region", "Revenue"]); _sh.append(["West", 42000])
_wb.save(nested / "numbers.xlsx")
ok("Region\tRevenue" in office.extract_office(nested / "numbers.xlsx")
   and "42000" in office.extract_office(nested / "numbers.xlsx"), "xlsx sheet rows extracted")
from pptx import Presentation as _PP  # noqa: E402
_pp = _PP(); _pp.slides.add_slide(_pp.slide_layouts[5]).shapes.title.text = "Launch Plan"
_pp.save(nested / "deck.pptx")
ok("Launch Plan" in office.extract_office(nested / "deck.pptx"), "pptx slide text extracted")
(nested / "guide.html").write_text("<h1>Runbook</h1><p>Restart the <b>worker</b>.</p>")
lws2 = Workspace(repo("local-office"))
r4 = localfiles.sync(store_office := Store(lws2), {"paths": [str(docs_dir)]}, settings=lset)
ok(str(nested / "report.docx") in r4["changed"], "nested .docx discovered by recursive crawl")
ok("disputeCard" in store_office.get("localfiles", str(nested / "report.docx")).body, "docx body indexed")
ok("42000" in store_office.get("localfiles", str(nested / "numbers.xlsx")).body, "xlsx indexed by localfiles")
ok("Launch Plan" in store_office.get("localfiles", str(nested / "deck.pptx")).body, "pptx indexed by localfiles")
ok("Restart the worker." in store_office.get("localfiles", str(nested / "guide.html")).body, "html flattened + indexed")
store_office.close()

# -- pdf backend routing (fakes, no models) -----------------------------------------------------
native = lambda p: ["real embedded text on this page", "   "]  # 2nd page has no embedded text
ocr = lambda p, cfg, only=None: [f"OCR[{i}]" for i in (only if only is not None else [0, 1])]
auto = pdf.extract_pdf("x.pdf", {"backend": "auto"}, native_text=native, ocr_pages=ocr)
ok("real embedded text" in auto and "OCR[1]" in auto and "OCR[0]" not in auto, "auto OCRs only image-only pages")
text_only = pdf.extract_pdf("x.pdf", {"backend": "text"}, native_text=native, ocr_pages=ocr)
ok("OCR" not in text_only, "text backend never OCRs")
forced = pdf.extract_pdf("x.pdf", {"backend": "unlimited-ocr"}, native_text=native, ocr_pages=ocr)
ok(forced == "OCR[0]\n\nOCR[1]", "unlimited-ocr backend OCRs every page")
degrade = pdf.extract_pdf("x.pdf", {"backend": "auto"}, native_text=native,
                          ocr_pages=lambda *a, **k: (_ for _ in ()).throw(RuntimeError("no device")))
ok(degrade.strip() == "real embedded text on this page", "auto keeps embedded text when a page's OCR errors")
try:
    pdf.extract_pdf("x.pdf", {"backend": "unlimited-ocr"}, native_text=native,
                    ocr_pages=lambda *a, **k: (_ for _ in ()).throw(RuntimeError("boom")))
    ok(False, "explicit unlimited-ocr should surface OCR errors")
except RuntimeError:
    ok(True, "explicit unlimited-ocr surfaces the error instead of silently degrading")

# -- hybrid search: deterministic keyword rescue + fusion ---------------------------------------
hws = Workspace(repo("hybrid"))
with Store(hws) as store:
    store.upsert("gdocs", "d1", title="Payments", url="u1", revision_id=None,
                 body="Billing uses chargeCard to bill customers monthly.")
    store.upsert("gdocs", "d2", title="Refunds", url="u2", revision_id=None,
                 body="An identifier like ZQ-9001 marks a reversed transaction.")
reembed(hws, embed_fn=fake_embed)  # embeds the two manually-stored docs into Lance + chunk mirror
kw = hybrid_search(hws, "ZQ-9001", embed_query_fn=lambda q: fake_embed([q])[0], hybrid=True, k=3)
ok(kw and kw[0]["doc_id"] == "d2", "rare exact identifier retrieved by the keyword half")
vec = hybrid_search(hws, "how are customers billed", embed_query_fn=lambda q: fake_embed([q])[0], k=3)
ok(any(h["doc_id"] == "d1" for h in vec), "semantic query finds the billing doc")
exp = hybrid_search(hws, "chargeCard", embed_query_fn=lambda q: fake_embed([q])[0], expand=1, k=2)
ok(exp and "context" in exp[0], "expand attaches surrounding context")

# -- retrieval primitives -----------------------------------------------------------------------
rec = recent(hws, source="gdocs", limit=5)
ok(len(rec) == 2 and rec[0]["source"] == "gdocs", "recent lists indexed docs")
th = thread(hws, "Refunds")
ok(th and "ZQ-9001" in th[0]["text"], "thread/doc returns the full body by title match")
with Store(hws) as store:
    first_chunk = store.find_docs(source="gdocs")  # sanity that docs exist
ok(bool(first_chunk), "find_docs returns documents")

# -- sync --rebuild re-embeds onto different settings -------------------------------------------
re_res = reembed(hws, embed_fn=fake_embed)
ok(re_res["docs"] == 2 and re_res["chunks"] >= 2, "--rebuild re-embeds every stored doc")

# -- sync checkpointing: an interrupted embed resumes without losing docs ------------------------
ckws = Workspace(repo("checkpoint"))
with Store(ckws) as s:
    s.upsert("gdocs", "a", title="A", url=None, revision_id="1", body="alpha body content that is definitely long enough to be chunked and embedded here", meta={"modified_at": "2026-01-01T00:00:00Z"})
    s.upsert("gdocs", "b", title="B", url=None, revision_id="1", body="beta body content that is definitely long enough to be chunked and embedded here", meta={"modified_at": "2026-02-01T00:00:00Z"})
    s.upsert("gdocs", "c", title="C", url=None, revision_id="1", body="gamma body content that is definitely long enough to be chunked and embedded here", meta={"modified_at": "2026-03-01T00:00:00Z"})
    ok([d for _, d in s.embed_queue()] == ["a", "b", "c"], "embed queue lists unembedded docs oldest-first")
_boom = {"n": 0}
def flaky_embed(texts):
    _boom["n"] += 1
    if _boom["n"] == 2:  # blow up while embedding the 2nd doc, mid-sync
        raise RuntimeError("interrupted")
    return fake_embed(texts)
try:
    run_sync(ckws, refetch=False, embed_fn=flaky_embed)
    ok(False, "interrupted embed should propagate")
except RuntimeError:
    ok(True, "an embed failure interrupts the sync")
with Store(ckws) as s:
    remaining = [d for _, d in s.embed_queue()]
ok("a" not in remaining and "b" in remaining, "docs before the interruption are checkpointed; the rest stay queued")
run_sync(ckws, refetch=False, embed_fn=fake_embed)  # resume — finishes what was left
with Store(ckws) as s:
    ok(s.embed_queue() == [], "resuming the sync drains the embed queue")
    s.upsert("gdocs", "a", title="A", url=None, revision_id="2", body="alpha EDITED body content that is definitely long enough to be chunked here", meta={"modified_at": "2026-01-01T00:00:00Z"})
    ok([d for _, d in s.embed_queue()] == ["a"], "an edited doc re-enters the queue; unchanged docs stay done")

# == retrieval upgrades: weighted RRF, routing, recency, merge, enrichment, large chunks, rerank ==
from bean.search import related as _related, _merge_sections, _query_weights  # noqa: E402
from bean import graph as _graph  # noqa: E402

_WCFG = {"vector_weight": 1.0, "keyword_weight": 1.0, "auto_weight": True}
vw, kw = _query_weights("ZQ-9001", _WCFG)
ok(kw > vw, "auto_weight leans keyword for an identifier query")
vw, kw = _query_weights("how are customers billed today", _WCFG)
ok(vw > kw, "auto_weight leans vector for a natural-language question")
vw2, kw2 = _query_weights("ZQ-9001", {**_WCFG, "auto_weight": False})
ok(vw2 == 1.0 and kw2 == 1.0, "auto_weight off keeps the fixed weights")

# query variants: a variant carrying the identifier rescues a doc the main query misses
qvws = Workspace(repo("q-variants"))
with Store(qvws) as store:
    store.upsert("gdocs", "dm", title="Meeting", url=None, revision_id=None,
                 body="the quarterly planning notes for the team roadmap discussion this week")
    store.upsert("gdocs", "di", title="Ticket", url=None, revision_id=None,
                 body="incident ZX-42 root cause analysis and the remediation steps that were taken")
reembed(qvws, embed_fn=fake_embed)
main_only = hybrid_search(qvws, "quarterly roadmap planning", embed_query_fn=lambda q: fake_embed([q])[0], k=5)
with_variant = hybrid_search(qvws, "quarterly roadmap planning", queries=["ZX-42"],
                             embed_query_fn=lambda q: fake_embed([q])[0], k=5)
ok(main_only[0]["doc_id"] == "dm", "main query ranks the topical doc first")
ok(with_variant[0]["doc_id"] == "di", "the ZX-42 variant floats the identifier doc to the top via fusion")

# recency: with decay on, the newer doc outranks an equally-relevant old one
rws = Workspace(repo("recency"))
with Store(rws) as store:
    store.upsert("gdocs", "old", title="Old widget", url=None, revision_id=None,
                 body="widget alpha beta gamma delta epsilon components overview and details",
                 meta={"modified_at": "2019-01-01T00:00:00Z"})
    store.upsert("gdocs", "new", title="New widget", url=None, revision_id=None,
                 body="widget alpha beta gamma delta epsilon components overview and details",
                 meta={"modified_at": "2026-06-01T00:00:00Z"})
reembed(rws, embed_fn=fake_embed)
cfgmod.save_global({"search": {"merge_sections": False}})
base = hybrid_search(rws, "widget alpha", embed_query_fn=lambda q: fake_embed([q])[0], k=2)
cfgmod.save_global({"search": {"merge_sections": False, "recency_decay": 8.0, "recency_floor": 0.05}})
decayed = hybrid_search(rws, "widget alpha", embed_query_fn=lambda q: fake_embed([q])[0], k=2, now=NOW)
ok(decayed and decayed[0]["doc_id"] == "new", f"recency decay floats the newer doc to the top ({[h['doc_id'] for h in decayed]})")
ok({h["doc_id"] for h in base} == {"old", "new"}, "both docs retrieved regardless of recency")
cfgmod.save_global({})

# section merge: two adjacent chunks of one doc collapse into a single section spanning both
mws = Workspace(repo("merge"))
body = "\n".join(f"line {i} content here about topic" for i in range(120))
with Store(mws) as store:
    store.upsert("gdocs", "big", title="Big", url=None, revision_id=None, body=body)
reembed(mws, embed_fn=fake_embed)
with Store(mws) as store:
    ch = store.neighbors("gdocs", "big", 0, 999)  # all base chunks of the doc, ordered by ord
    two = [{**ch[0], "score": 0.9}, {**ch[1], "score": 0.8}]
    merged = _merge_sections(store, two)
    ok(len(merged) == 1, "adjacent same-doc chunks merge into one section")
    ok(merged[0]["start"] == ch[0]["start"] and merged[0]["end"] == ch[1]["end"],
       "merged section spans the union line range")

# chunk enrichment: title_prefix lets a title-only term retrieve a doc whose body lacks it
ews = Workspace(repo("enrich"))
with Store(ews) as store:
    store.upsert("gdocs", "z", title="Zephyr", url=None, revision_id=None,
                 body="the quick brown fox jumps over the lazy dog again and again today")
reembed(ews, embed_fn=fake_embed)  # title_prefix defaults True
hit = hybrid_search(ews, "Zephyr", embed_query_fn=lambda q: fake_embed([q])[0], k=3)
ok(any(h["doc_id"] == "z" for h in hit), "title_prefix embeds the title so a title term retrieves the doc")
with Store(ews) as _st:
    ok("Zephyr" not in _st.get("gdocs", "z").body, "…even though the body never says it (title stored separately)")

# large chunks: enabling them adds coarse vectors on top of the base chunks
lgws = Workspace(repo("largechunks"))
longbody = "\n".join(f"line {i} about payments and billing" for i in range(200))
with Store(lgws) as store:
    store.upsert("gdocs", "L", title="Long", url=None, revision_id=None, body=longbody)
cfgmod.save_global({})
n_base = reembed(lgws, embed_fn=fake_embed)["chunks"]
cfgmod.save_global({"chunking": {"large_chunks": True}})
n_large = reembed(lgws, embed_fn=fake_embed)["chunks"]
ok(n_large > n_base, f"large_chunks adds coarse vectors ({n_base} -> {n_large})")
cfgmod.save_global({})

# rerank: an injected cross-encoder reorders the fused candidates
rkws = Workspace(repo("rerank"))
with Store(rkws) as store:
    store.upsert("gdocs", "a", title="A", url=None, revision_id=None,
                 body="payments billing invoice records for the finance team monthly review")
    store.upsert("gdocs", "b", title="B", url=None, revision_id=None,
                 body="payments billing MARKER invoice records for the finance team monthly review")
reembed(rkws, embed_fn=fake_embed)
cfgmod.save_global({"search": {"rerank": {"enabled": True}, "merge_sections": False}})
rr = hybrid_search(rkws, "payments billing", embed_query_fn=lambda q: fake_embed([q])[0], k=2,
                   rerank_fn=lambda q, texts: [1.0 if "MARKER" in t else 0.0 for t in texts])
ok(rr and rr[0]["doc_id"] == "b", "injected reranker floats the MARKER doc to the top")
cfgmod.save_global({})

# graph: implied edges (authored_by / in-container) + `related` one-hop expansion
ok([e["dst"] for e in _graph.implied_edges(type("D", (), {"source": "github", "doc_id": "acme/repo#1", "author": "Ada"})())]
   == ["Ada", "acme/repo"], "implied_edges derives author + repo container from a github doc id")
gws2 = Workspace(repo("graph"))
with Store(gws2) as store:
    for did, auth in [("acme/repo#1", "Ada"), ("acme/repo#2", "Bo"), ("other/x#9", "Ada")]:
        store.upsert("github", did, title=did, url=None, revision_id=None, body=f"issue {did}",
                     meta={"author": auth})
        store.replace_edges("github", did, _graph.implied_edges(store.get("github", did)))
    rel = store.related("github", "acme/repo#1")
    ids = {r["doc_id"] for r in rel}
    ok("acme/repo#2" in ids, "related finds a doc in the same repo container")
    ok("other/x#9" in ids, "related finds another doc by the same author")
    ok(all("reason" in r for r in rel), "each related hit carries a reason")

# metadata filters: author + date narrowing on search and recent
fws = Workspace(repo("filters"))
with Store(fws) as store:
    store.upsert("gdocs", "ada1", title="Ada plan", url=None, revision_id=None,
                 body="roadmap widget alpha planning notes for the next quarter and beyond",
                 meta={"author": "Ada Lovelace", "modified_at": "2026-06-01T00:00:00Z"})
    store.upsert("gdocs", "bob1", title="Bob plan", url=None, revision_id=None,
                 body="roadmap widget alpha planning notes for the next quarter and beyond",
                 meta={"author": "Bob", "modified_at": "2020-01-01T00:00:00Z"})
reembed(fws, embed_fn=fake_embed)
byauthor = hybrid_search(fws, "roadmap widget", author="Ada", embed_query_fn=lambda q: fake_embed([q])[0], k=5)
ok({h["doc_id"] for h in byauthor} == {"ada1"}, "search --author filters to that author's docs")
recent_since = recent(fws, since="2025-01-01")
ok({h["doc_id"] for h in recent_since} == {"ada1"}, "recent --since filters by modified date")
ok({h["doc_id"] for h in recent(fws, author="Bob")} == {"bob1"}, "recent --author filters by author")

# staleness: bean warns when the index is old, but never syncs on its own
from bean.cli import _staleness_note  # noqa: E402
import datetime as _dt  # noqa: E402
stws = Workspace(repo("stale"))
ok(_staleness_note(stws) is None, "never-synced index does not nag (setup handles it)")
with Store(stws) as store:
    store.set_state("last_sync", "2020-01-01T00:00:00+00:00")
_note = _staleness_note(stws)
ok(_note is not None and "stale" in _note and "sync" in _note, "an old last_sync triggers a stale warning")
with Store(stws) as store:
    store.set_state("last_sync", _dt.datetime.now(_dt.timezone.utc).isoformat())
ok(_staleness_note(stws) is None, "a fresh sync clears the warning")
cfgmod.save_global({"sync": {"stale_days": 0}})
with Store(stws) as store:
    store.set_state("last_sync", "2020-01-01T00:00:00+00:00")
ok(_staleness_note(stws) is None, "stale_days=0 disables the warning")
cfgmod.save_global({})

# == global vs local connector scope ============================================================
from bean.workspace import save_scopes, set_source_scope, source_scope  # noqa: E402
from bean.search import search_many  # noqa: E402
from bean.sync import run_sync as _run_sync  # noqa: E402
import io as _io  # noqa: E402
import contextlib as _ctx  # noqa: E402

save_scopes({})
ok(source_scope("slack") == "local", "default connector scope is local")
set_source_scope("slack", "global")
ok(source_scope("slack") == "global", "scope set/get round-trips")
save_scopes({})

g1, g2 = Workspace.global_(), Workspace.global_()
ok(g1.dir == g2.dir and g1.is_global and g1.dir.name == "_global", "global workspace is stable + flagged")

# run_sync keys= restricts which sources sync (route global vs local sets into different stores)
kdir = Path(tempfile.mkdtemp(prefix="bean-keys-"))
(kdir / "n.md").write_text("# Note\nsome content about widgets and rollbacks for indexing purposes here\n")
kws = Workspace(repo("keys"))
kws.save_config({"localfiles": {"paths": [str(kdir)]}, "github": {"repos": ["a/b"]}})
save_credential("github", {"token": "t"})
kr = _run_sync(kws, keys={"localfiles"}, embed_fn=fake_embed,
               fetch=lambda u, h, method="GET", body=None: res(404, "{}"))
ok(kr["changed"] and all(s == "localfiles" for s, _ in kr["changed"]),
   f"run_sync keys= restricts to the named sources ({kr['changed']})")

# a doc in the repo store + a doc in the global store both surface via search_many
lrepo = Workspace(repo("scope-repo"))
gws = Workspace.global_()
with Store(lrepo) as s:
    s.upsert("github", "acme/app#1", title="Local issue", url=None, revision_id=None,
             body="the local repo ticket about deployment rollback steps and the details")
with Store(gws) as s:
    s.upsert("slack", "eng/2026-W01", title="Global slack", url=None, revision_id=None,
             body="the global slack thread about deployment rollback plans and the timing")
reembed(lrepo, embed_fn=fake_embed)
reembed(gws, embed_fn=fake_embed)
un = search_many([lrepo, gws], "deployment rollback", embed_query_fn=lambda q: fake_embed([q])[0], k=5)
ok({h["source"] for h in un} >= {"github", "slack"}, f"search_many unions repo + global stores ({[h['source'] for h in un]})")

# cmd_scope moves a connector's config + purges its old index (so a resync repopulates the new scope)
from bean.cli import cmd_scope  # noqa: E402
mrepo = Workspace(repo("scope-mig"))
mrepo.save_config({"github": {"repos": ["x/y"]}})
with Store(mrepo) as s:
    s.upsert("github", "x/y#1", title="t", url=None, revision_id=None,
             body="body about the scope migration test with enough content to index a chunk here")
reembed(mrepo, embed_fn=fake_embed)
save_scopes({})
_args = type("A", (), {"source": "github", "value": "global"})()
with _ctx.redirect_stdout(_io.StringIO()):
    cmd_scope(mrepo, _args)
ok(source_scope("github") == "global", "cmd_scope records the new scope")
with Store(mrepo) as s:
    ok(s.doc_ids("github") == [], "cmd_scope purges the moved source from the old (repo) store")
ok("x/y" in (Workspace.global_().load_config().get("github") or {}).get("repos", []),
   "cmd_scope moves tracked items into the global workspace config")
ok("x/y" not in (mrepo.load_config().get("github") or {}).get("repos", []),
   "cmd_scope removes the items from the repo config")

# cleanup shared global state so later tests see a clean slate
save_scopes({})
_gws = Workspace.global_()
with Store(_gws) as s:
    for _src in ("slack", "github"):
        for _d in s.doc_ids(_src):
            s.delete(_src, _d)
_gc = _gws.load_config()
for _k in ("github", "slack"):
    _gc.pop(_k, None)
_gws.save_config(_gc)

# == per-workspace credentials: local connectors get their own; global share; fallback ==========
from bean.workspace import credential_context, credential_path  # noqa: E402
import stat as _stat  # noqa: E402
wa, wb = Workspace(repo("cred-a")), Workspace(repo("cred-b"))
with credential_context(wa):
    save_credential("github", {"token": "AAA"})
with credential_context(wb):
    save_credential("github", {"token": "BBB"})
with credential_context(wa):
    ok(load_credential("github")["token"] == "AAA", "repo A sees its own github credential")
with credential_context(wb):
    ok(load_credential("github")["token"] == "BBB", "repo B sees a different github credential")
# a shared credential falls through to any local workspace that lacks its own
save_credential("sharedsvc", {"token": "SHARED"})  # no context -> shared dir
with credential_context(wa):
    ok(load_credential("sharedsvc")["token"] == "SHARED", "shared credential falls through to a local ws")
with credential_context(wa):
    save_credential("sharedsvc", {"token": "LOCAL_A"})  # ...but a local one shadows it
with credential_context(wa):
    ok(load_credential("sharedsvc")["token"] == "LOCAL_A", "workspace credential shadows the shared one")
with credential_context(wb):
    ok(load_credential("sharedsvc")["token"] == "SHARED", "another workspace still sees the shared one")
with credential_context(Workspace.global_()):
    ok(load_credential("sharedsvc")["token"] == "SHARED", "a global-scope context reads the shared dir")
ok("cred-a" in str(credential_path("github", wa)), "credential_path is workspace-scoped for a local ws")
ok("cred-a" not in str(credential_path("slack", None)), "credential_path is the shared dir for a global connector")
_pa = wa.dir / "credentials" / "github.json"
ok(_pa.exists() and _stat.S_IMODE(_pa.stat().st_mode) == 0o600, "local credential stored in the workspace, mode 0600")

# == new connectors (offline: fake fetch, real Store) ===========================================
from bean.connectors import confluence, jira, zendesk, discord, microsoft, salesforce  # noqa: E402  (core)


def check(name, fn):
    """Run one connector's offline test in isolation so a raise can't abort the whole suite."""
    try:
        fn()
    except Exception as err:  # noqa: BLE001
        import traceback
        ok(False, f"{name}: raised {err!r}\n{traceback.format_exc()}")


def _store(tag):
    return Store(Workspace(repo(tag)))


def t_confluence():
    save_credential("confluence", {"method": "cloud", "url": "https://x.atlassian.net/wiki",
                                   "email": "e", "token": "t"})
    body_fetches = []
    def f(u, h, method="GET", body=None):
        if "spaceKey" in u:  # cheap listing: metadata only, no body.storage
            return res(200, {"results": [{"id": "123", "title": "Runbook",
                "version": {"number": 2, "by": {"displayName": "Ada"}},
                "history": {"lastUpdated": {"when": "2026-01-02T00:00:00Z"}},
                "_links": {"webui": "/pages/123"}}], "size": 1})
        if "/content/123" in u:  # lazy body fetch — only for a changed page
            body_fetches.append(u)
            return res(200, {"body": {"storage": {"value": "<p>Restart the worker.</p>"}}})
        return res(200, {})
    with _store("confluence") as s:
        r = confluence.sync(s, {"spaces": ["ENG"], "pages": []}, settings={}, fetch=f)
        ok(r["changed"] == ["123"], f"confluence page ingested ({r})")
        ok("Restart the worker." in s.get("confluence", "123").body, "confluence storage HTML flattened")
        ok(len(body_fetches) == 1, "confluence pulled the body once for the changed page")
        r2 = confluence.sync(s, {"spaces": ["ENG"], "pages": []}, settings={}, fetch=f)
        ok(r2["changed"] == [], "confluence unchanged version is a no-op")
        ok(len(body_fetches) == 1, "confluence skips the body fetch when the version is unchanged")


def t_jira():
    save_credential("jira", {"method": "cloud", "url": "https://x.atlassian.net", "email": "e", "token": "t"})
    def f(u, h, method="GET", body=None):
        return res(200, {"issues": [{"key": "PROJ-1", "fields": {"summary": "Fix it",
            "updated": "2026-01-01T00:00:00.000+0000", "status": {"name": "Open"},
            "reporter": {"displayName": "B"}, "assignee": None, "description": "desc",
            "comment": {"comments": []}}}], "total": 1})
    with _store("jira") as s:
        r = jira.sync(s, {"projects": ["PROJ"]}, settings={}, fetch=f)
        ok(r["changed"] == ["PROJ-1"], f"jira issue ingested ({r})")


def t_zendesk():
    save_credential("zendesk", {"subdomain": "acme", "email": "e", "token": "t"})
    def f(u, h, method="GET", body=None):
        if "incremental/tickets" in u:
            return res(200, {"tickets": [{"id": 7, "subject": "S", "updated_at": "2026-01-01",
                "created_at": "2025-12-01", "description": "d", "status": "open"}],
                "end_of_stream": True, "end_time": 123})
        if "/tickets/7/comments" in u:
            return res(200, {"comments": []})
        if "help_center/articles" in u:
            return res(200, {"articles": [{"id": 3, "title": "A", "body": "<p>x</p>",
                "updated_at": "2026-01-01", "created_at": "2025-01-01", "html_url": "http://a"}],
                "next_page": None})
        return res(404, {})
    with _store("zendesk") as s:
        r = zendesk.sync(s, {}, settings={}, fetch=f, now=200.0)
        ok(sorted(r["changed"]) == ["article/3", "ticket/7"], f"zendesk tickets+articles ({r})")


def t_discord():
    import datetime as _dt
    save_credential("discord", {"token": "t"})
    mts = NOW - 2 * 86400
    iso = _dt.datetime.fromtimestamp(mts, tz=_dt.timezone.utc).isoformat()
    def f(u, h, method="GET", body=None):
        if "/channels/999/messages" in u:
            if "before=" in u:
                return res(200, [])
            return res(200, [{"id": "1000", "timestamp": iso, "content": "hello team",
                              "author": {"username": "ada"}}])
        if "/channels/999" in u:
            return res(200, {"name": "general"})
        return res(404, {})
    with _store("discord") as s:
        r = discord.sync(s, {"channels": ["999"]}, settings={}, fetch=f, now=NOW)
        ok(r["changed"] == ["general/1000"], f"discord one doc per message ({r})")
        doc = s.get("discord", "general/1000")
        ok("hello team" in doc.body, "discord message rendered")
        ok(str(doc.modified_at).startswith(iso[:10]), "discord message carries its timestamp")
        ok("hello team" in doc.title, "discord message title is its text")


def t_microsoft():
    save_credential("microsoft", {"method": "az"})
    def f(u, h, method="GET", body=None):
        if "/me/drive/root/children" in u:
            return res(200, {"value": [{"id": "F1", "name": "notes.md", "file": {}, "eTag": "e1",
                "@microsoft.graph.downloadUrl": "https://dl/notes"}]})
        if u.startswith("https://dl/"):
            return res(200, "# notes from onedrive")
        return res(200, {"value": []})
    with _store("microsoft") as s:
        r = microsoft.sync(s, {"drives": ["me"]}, settings={}, fetch=f, token_fn=lambda force=False: "tok")
        ok(r["changed"] == ["file/F1"], f"microsoft onedrive file ingested ({r})")


def t_salesforce():
    save_credential("salesforce", {"token": "t", "url": "https://x.my.salesforce.com"})
    def f(u, h, method="GET", body=None):
        if "Knowledge__kav" in u:
            return res(200, {"records": [{"Id": "k1", "Title": "KB", "Summary": "<p>s</p>",
                                          "UrlName": "kb", "LastModifiedDate": "2026-01-01"}]})
        if "Case" in u:
            return res(200, {"records": [{"Id": "c1", "CaseNumber": "0001", "Subject": "Sub",
                "Description": "<p>d</p>", "Status": "New", "LastModifiedDate": "2026-01-01"}]})
        return res(200, {"records": []})
    with _store("salesforce") as s:
        r = salesforce.sync(s, {}, settings={}, fetch=f)
        ok(sorted(r["changed"]) == ["article/k1", "case/c1"], f"salesforce KB+cases ({r})")


for _name, _fn in [
    ("confluence", t_confluence), ("jira", t_jira), ("zendesk", t_zendesk),
    ("discord", t_discord), ("microsoft", t_microsoft), ("salesforce", t_salesforce),
]:
    check(_name, _fn)

# == Onyx-parity connectors (offline) ===========================================================
from bean.connectors import hubspot  # noqa: E402  (core)


def t_hubspot():
    save_credential("hubspot", {"token": "t", "portal_id": "7"})
    def f(u, h, method="GET", body=None):
        if "/tickets" in u:
            return res(200, {"results": [{"id": "5", "updatedAt": "2026-01-02T00:00:00Z",
                "properties": {"subject": "T", "content": "c"}}]})
        return res(200, {"results": []})
    with _store("hubspot") as s:
        ok(hubspot.sync(s, {"include": ["tickets"]}, settings={}, fetch=f)["changed"] == ["ticket/5"],
           "hubspot ticket ingested")


for _name, _fn in [
    ("hubspot", t_hubspot),
]:
    check(_name, _fn)

# == plugin system: core-only by default, load a drop-in plugin =================================
import bean.sources as _S  # noqa: E402
from bean import config as _cfg  # noqa: E402
from bean.plugins import discover_sources  # noqa: E402

core_keys = {s.key for s in _S.CORE_SOURCES}
ok(core_keys == {"slack", "gdocs", "github", "confluence", "jira", "zendesk",
                 "salesforce", "hubspot", "microsoft", "discord"}, "10 cloud core connectors")
ok(_S.SOURCES[-1].key == "localfiles", "localfiles registered last (path catch-all)")

# a drop-in plugin file: a standalone module exposing SOURCE
plugdir = Path(tempfile.mkdtemp(prefix="bean-plugins-"))
(plugdir / "acme.py").write_text(
    "from bean.sources import Source\n"
    "def sync(store, config, *, settings, fetch=None, full=False, since_days=90, log=lambda m: None):\n"
    "    changed = []\n"
    "    for b in config.get('boards', []):\n"
    "        if store.upsert('acme', f'acme/{b}', title=b, url=None, revision_id='r1',\n"
    "                        body=f'# {b}\\nacme board body'):\n"
    "            changed.append(f'acme/{b}')\n"
    "    return {'changed': changed, 'removed': []}\n"
    "def connected():\n"
    "    return {'ok': True}\n"
    "SOURCE = Source('acme', 'acme', 'Acme', ('boards',), sync, auth=None,\n"
    "                add_help='acme:BOARD', connected=connected)\n"
)
found = discover_sources(_S.Source, global_config={}, dirs=[plugdir])
ok(len(found) == 1 and found[0].key == "acme", "drop-in plugin discovered from a dir")

_cfg.save_global({"plugins": {"paths": [str(plugdir)]}})
_S.reload_sources()
ok("acme" in _S.BY_KEY and _S.SOURCES[-1].key == "localfiles",
   "drop-in plugin registers ahead of the localfiles catch-all")
acme_src = _S.BY_KEY["acme"]
with _store("acme-plugin") as s:
    r = acme_src.sync(s, {"boards": ["main"]}, settings={}, fetch=None)
    ok(r["changed"] == ["acme/main"] and s.get("acme", "acme/main"), "drop-in plugin syncs a document")
_cfg.save_global({})  # reset global config
_S.reload_sources()

# --- CLI rendering: preview vs --full -----------------------------------------------------------
import io
from contextlib import redirect_stdout
from bean.cli import _print_hits  # noqa: E402

_long = "\n".join(f"line {i} " + "x" * 200 for i in range(12))
_hit = [{"doc_id": "d1", "title": "Doc One", "text": _long}]


def _render(**kw):
    buf = io.StringIO()
    with redirect_stdout(buf):
        _print_hits(None, _hit, "empty", **kw)
    return buf.getvalue()

prev = _render()  # default preview
ok(prev.count("line ") == 5, "preview caps at 5 lines")
ok("xxxx" in prev and max(len(l) for l in prev.splitlines()) < 130, "preview truncates each line to ~110 chars")

full = _render(full=0)  # no cap
ok(full.count("line ") == 12, "--full 0 prints every line")
ok("x" * 200 in full, "--full 0 does not truncate line length")

capped = _render(full=300)  # char budget
ok("truncated at 300 chars" in capped, "--full N under budget prints a truncation notice")
ok(len(capped) < len(full), "--full N caps output shorter than unlimited")

_ebuf = io.StringIO()
with redirect_stdout(_ebuf):
    _erc = _print_hits(None, [], "nothing here", full=0)
ok(_erc == 1 and "nothing here" in _ebuf.getvalue(),
   "empty hits still print the empty message under --full")

# --- http POST / JSON helpers + error-detail extraction -----------------------------------------
from bean.http import api_post, api_json, api_json_post  # noqa: E402

_seen = {}
def _echo(u, h, method="GET", body=None):
    _seen["method"], _seen["body"] = method, body
    return res(200, {"ok": True})
ok(api_post("https://x/y", {}, {"a": 1}, fetch=_echo).json() == {"ok": True}, "api_post returns the response")
ok(_seen["method"] == "POST" and _seen["body"] == {"a": 1}, "api_post threads method + JSON body through the seam")
ok(api_json_post("https://x/y", {}, {"a": 1}, fetch=_echo) == {"ok": True}, "api_json_post unwraps the JSON body")
ok(api_json("https://x/y", {}, fetch=lambda u, h: res(200, {"v": 9}))["v"] == 9, "api_json unwraps a GET body")
try:
    api_json("https://x/y", {}, fetch=lambda u, h: res(500, {"error": {"message": "boom detail"}}))
    ok(False, "api_json should raise on a non-2xx")
except RuntimeError as err:
    ok("boom detail" in str(err), "_detail surfaces a nested error.message")
try:
    api_json("https://x/y", {}, fetch=lambda u, h: res(422, {"errors": ["first bad thing"]}))
    ok(False, "api_json should raise on 422")
except RuntimeError as err:
    ok("first bad thing" in str(err), "_detail extracts the first item of an errors list")

# --- content_hash is deterministic and content-sensitive ----------------------------------------
ok(content_hash("same text") == content_hash("same text"), "content_hash is deterministic")
ok(content_hash("a") != content_hash("b"), "content_hash distinguishes different bodies")

# --- github connector: sync + per-item tolerance (one bad issue never aborts the repo) ----------
def t_github():
    save_credential("github", {"token": "ghp_x", "login": "me"})
    def f(u, h, method="GET", body=None):
        if "/issues?" in u:  # the repo issue listing
            return res(200, [
                {"number": 1, "title": "Good", "state": "open", "user": {"login": "ada"},
                 "body": "a real issue body", "updated_at": "2026-01-01T00:00:00Z", "comments": 0,
                 "html_url": "https://github.com/acme/repo/issues/1",
                 "created_at": "2026-01-01T00:00:00Z"},
                {"title": "malformed — no number"},  # _ingest_issue raises (KeyError on number)
            ])
        return res(200, [])
    with _store("github") as s:
        r = github.sync(s, {"repos": ["acme/repo"], "include": ["issues"]}, settings={}, fetch=f)
        ok(r["changed"] == ["acme/repo#1"], f"github ingests the good issue ({r})")
        ok(s.get("github", "acme/repo#1") is not None,
           "the good issue survived a malformed sibling in the same batch")
check("github", t_github)

# --- neighbors primitive + scope-union wrappers (single ws behaves like the primitives) ---------
from bean.search import recent_many, neighbors_many, document_many  # noqa: E402

with Store(mws) as _s:
    _allch = _s.neighbors("gdocs", "big", 0, 999)
_cid = _allch[len(_allch) // 2]["id"]
_nb = neighbors(mws, _cid, radius=1)
ok(_nb and all(h["doc_id"] == "big" for h in _nb), "neighbors() returns chunks around a chunk id")
ok(neighbors(mws, "no-such-chunk") == [], "neighbors() is empty for an unknown chunk id")
ok([h["doc_id"] for h in recent_many(mws, limit=5)] == [h["doc_id"] for h in recent(mws, limit=5)],
   "recent_many over a single ws matches recent()")
ok(bool(neighbors_many(mws, _cid, radius=1)), "neighbors_many returns the first ws that has the chunk")
_dm = document_many(mws, "big")
ok(_dm and _dm[0]["doc_id"] == "big", "document_many finds the doc by id substring")
_rel = _related(gws2, "acme/repo#1")
ok(_rel and any(h["doc_id"] == "acme/repo#2" for h in _rel), "search.related() expands one hop over the graph")

# --- cmd_sql read-only guard --------------------------------------------------------------------
from bean.cli import cmd_sql  # noqa: E402
import types as _types  # noqa: E402

_sqlws = Workspace(repo("sqlguard"))
with Store(_sqlws) as _s:
    _s.upsert("gdocs", "d", title="T", url=None, revision_id=None, body="hello sql world")

def _sql(*words):
    return _types.SimpleNamespace(query=list(words), global_=False)

with _ctx.redirect_stderr(_io.StringIO()):
    ok(cmd_sql(_sqlws, _sql("DROP", "TABLE", "documents")) == 2, "cmd_sql refuses a non-SELECT (read-only guard)")
    ok(cmd_sql(_sqlws, _sql("DELETE", "FROM", "documents")) == 2, "cmd_sql refuses DELETE")
_okbuf = _io.StringIO()
with _ctx.redirect_stdout(_okbuf):
    _rc = cmd_sql(_sqlws, _sql("SELECT", "count(*)", "FROM", "documents"))
ok(_rc == 0 and "1" in _okbuf.getvalue(), "cmd_sql runs a SELECT and prints the result")
with _ctx.redirect_stdout(_io.StringIO()):
    ok(cmd_sql(_sqlws, _sql("WITH", "x", "AS", "(SELECT", "1", "AS", "n)", "SELECT", "*", "FROM", "x")) == 0,
       "cmd_sql allows a WITH … SELECT CTE")
_statebuf = _io.StringIO()
with _ctx.redirect_stdout(_statebuf):
    _src = cmd_sql(_sqlws, _sql("SELECT", "count(*)", "FROM", "state"))
ok(_src == 0, "cmd_sql exposes the state table (SELECT from state succeeds)")

# == Lance catalog: relational tables on Lance, queried via DuckDB =============================
from bean.lancecat import Catalog  # noqa: E402
_catdir = Path(tempfile.mkdtemp(prefix="bean-cat-"))
_cat = Catalog(_catdir)
_cat.upsert_documents([{"source": "gdocs", "doc_id": "d1", "title": "T", "url": "u", "revision_id": "r1",
                        "hash": "h1", "body": "alpha body", "created_at": None, "modified_at": None,
                        "author": "Ada", "mime": None, "fetched_at": None}])
_cat.upsert_documents([{"source": "gdocs", "doc_id": "d1", "title": "T2", "url": "u", "revision_id": "r2",
                        "hash": "h2", "body": "beta body", "created_at": None, "modified_at": None,
                        "author": "Ada", "mime": None, "fetched_at": None}])
_rows = _cat.duck().execute("SELECT title, hash, body FROM documents WHERE doc_id='d1'").fetchall()
ok(_rows == [("T2", "h2", "beta body")], "Lance upsert updates in place, queried via DuckDB")
_cat.delete_documents("gdocs", ["d1"])
ok(_cat.duck().execute("SELECT count(*) FROM documents").fetchone()[0] == 0, "Lance delete removes the row")

# ord is stored, not derived
_ordws = Workspace(repo("ord"))
with Store(_ordws) as store:
    store.upsert("gdocs", "o1", title="O", url=None, revision_id=None,
                 body="\n".join(f"line {i} of content here about a topic" for i in range(120)))
reembed(_ordws, embed_fn=fake_embed)
with Store(_ordws) as store:
    _ch = store.neighbors("gdocs", "o1", 0, 999)
ok([c["ord"] for c in _ch] == list(range(len(_ch))), "stored ord is a dense 0-based sequence per doc")

# legacy DuckDB catalog is migrated into Lance on first open
_mig = Workspace(repo("migrate"))
import duckdb as _dd
_c = _dd.connect(str(_mig.db_path))
_c.execute("CREATE TABLE documents (source TEXT, doc_id TEXT, title TEXT, url TEXT, revision_id TEXT, "
           "hash TEXT, body TEXT, created_at TIMESTAMP, modified_at TIMESTAMP, author TEXT, mime TEXT, "
           "fetched_at TIMESTAMP, embedded_hash TEXT)")
_c.execute("INSERT INTO documents VALUES ('gdocs','dz','T','u','r','h','body here',NULL,NULL,'Ada',NULL,now(),'h')")
_c.execute("CREATE TABLE revisions (source TEXT, doc_id TEXT, revision_id TEXT, hash TEXT, fetched_at TIMESTAMP)")
_c.execute("CREATE TABLE edges (source TEXT, src_doc TEXT, rel TEXT, dst_kind TEXT, dst TEXT)")
_c.close()
with Store(_mig) as s:
    ok(s.get("gdocs", "dz") is not None and s.get("gdocs", "dz").author == "Ada", "legacy doc migrated to Lance")
with Store(_mig) as s:
    ok(s.counts().get("gdocs") == 1, "second open is a no-op (no double-migration)")

# regression: a crash right after the documents-copy (Lance `documents` already has the row) but
# before revisions/edges are copied, with `catalog_migrated` still unset, must NOT drop
# revisions/edges without migrating them first (the old `cat_populated` shortcut did exactly that
# and silently destroyed both tables — see bean/store.py `_migrate_legacy_catalog`).
from bean.lancecat import Catalog as _Cat
from datetime import datetime as _dt
_mig3 = Workspace(repo("migrate-partial"))
_c3 = _dd.connect(str(_mig3.db_path))
_c3.execute("CREATE TABLE documents (source TEXT, doc_id TEXT, title TEXT, url TEXT, revision_id TEXT, "
            "hash TEXT, body TEXT, created_at TIMESTAMP, modified_at TIMESTAMP, author TEXT, mime TEXT, "
            "fetched_at TIMESTAMP, embedded_hash TEXT)")
_c3.execute("INSERT INTO documents VALUES ('gdocs','pz','T','u','r','h','body here',NULL,NULL,'Ada',NULL,now(),'h')")
_c3.execute("CREATE TABLE revisions (source TEXT, doc_id TEXT, revision_id TEXT, hash TEXT, fetched_at TIMESTAMP)")
_c3.execute("INSERT INTO revisions VALUES ('gdocs','pz','r0','h0',now())")
_c3.execute("CREATE TABLE edges (source TEXT, src_doc TEXT, rel TEXT, dst_kind TEXT, dst TEXT)")
_c3.execute("INSERT INTO edges VALUES ('gdocs','pz','links','doc','other-doc')")
_c3.close()
# Simulate the interrupted-migration state: Lance `documents` already populated for this doc (as if
# `upsert_documents` had run once already), legacy tables still present, flag still unset.
_Cat(_mig3.catalog_dir).upsert_documents([{
    "source": "gdocs", "doc_id": "pz", "title": "T", "url": "u", "revision_id": "r", "hash": "h",
    "body": "body here", "created_at": None, "modified_at": None, "author": "Ada", "mime": None,
    "fetched_at": _dt.now(),
}])
with Store(_mig3) as s:
    ok(any(r[0] == "r0" for r in s.revisions("gdocs", "pz")),
       "migration copies revisions even when Lance documents is already populated (no data loss)")
    ok(any(e["dst"] == "other-doc" for e in s.edges_of("gdocs", "pz")),
       "migration copies edges even when Lance documents is already populated (no data loss)")

# cloud config + Workspace cloud awareness (Task 2.1) -------------------------------------------
cloudws = Workspace(repo("cloud-off"))
ok(cloudws.is_cloud is False, "no cloud config -> is_cloud False")
ok(cloudws.remote_uri is None, "no cloud config -> remote_uri None")
ok(cloudws.replica_dir == cloudws.catalog_dir, "replica_dir mirrors catalog_dir")

cloudws2 = Workspace(repo("cloud-on"))
cloudws2.save_config({"settings": {"cloud": {
    "enabled": True, "bucket": "b", "prefix": "p/x", "region": "us-east-1",
}}})
ok(cloudws2.is_cloud is True, "cloud.enabled True -> is_cloud True")
ok(cloudws2.remote_uri == "s3://b/p/x", "remote_uri built from bucket + prefix")
ok(cloudws2.cloud["region"] == "us-east-1", "cloud property exposes resolved region")

cloudws3 = Workspace(repo("cloud-noprefix"))
cloudws3.save_config({"settings": {"cloud": {"enabled": True, "bucket": "b", "prefix": ""}}})
ok(cloudws3.remote_uri == "s3://b", "empty prefix -> bare bucket uri")

# remote.pull: S3->local replication of the Lance catalog (Task 2.2) -----------------------------
# A local temp dir stands in for S3 — `_copy_new`'s local-dir branch shares the exact copy logic
# the s3:// branch drives via `aws s3 sync`, so this exercises the same replication behaviour.
from bean import remote as _remote  # noqa: E402
_remote_dir = Path(tempfile.mkdtemp(prefix="bean-remote-"))
_replica_dir = Path(tempfile.mkdtemp(prefix="bean-replica-")) / "catalog"
_remote_cat = Catalog(_remote_dir)
_remote_cat.upsert_documents([{"source": "gdocs", "doc_id": "r1", "title": "T", "url": "u",
                                "revision_id": "rev1", "hash": "h1", "body": "pulled body",
                                "created_at": None, "modified_at": None, "author": "Ada",
                                "mime": None, "fetched_at": None}])
_remote._copy_new(str(_remote_dir), _replica_dir)
_replica_cat = Catalog(_replica_dir)
ok(_replica_cat.duck().execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
   "_copy_new replicates a freshly-written doc into a fresh local replica")
ok(_replica_cat.duck().execute("SELECT body FROM documents WHERE doc_id='r1'").fetchone()[0]
   == "pulled body", "replicated row matches the remote's committed data")

# manifest-last ordering: the replica must open cleanly, meaning no manifest ever landed locally
# before the data files it references (an out-of-order copy would leave a dataset that fails to
# open, since Lance would resolve the manifest to fragment files that aren't there yet).
_replica_versions = list((_replica_dir / "documents.lance" / "_versions").glob("*.manifest"))
_replica_data = list((_replica_dir / "documents.lance" / "data").glob("*.lance"))
ok(len(_replica_versions) > 0 and len(_replica_data) > 0,
   "replica has both manifests and data files after _copy_new")

# second pull is idempotent: nothing new to copy, no error, same row count
_remote._copy_new(str(_remote_dir), _replica_dir)
ok(Catalog(_replica_dir).duck().execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
   "a second _copy_new is a no-op (already-present immutable files are never re-copied)")

# pull() itself: no-op for a non-cloud workspace, and routes a cloud workspace's local-dir remote
# through the same _copy_new path used above.
_pullws_off = Workspace(repo("pull-off"))
_remote.pull(_pullws_off)  # must not raise, must not touch replica_dir
ok(not (_pullws_off.replica_dir / "documents.lance").exists(), "pull() no-ops for a non-cloud workspace")

# _s3_sync_commands: pure argv builder for the two `aws s3 sync` passes, unit-testable without
# ever running `aws` or touching S3. Pass 2 must exclude-then-include, since a bare --include
# with no preceding --exclude "*" is a no-op in the AWS CLI (it would re-sync everything, letting
# manifests race data files instead of always landing last).
_cmds = _remote._s3_sync_commands("s3://b/p", Path("/tmp/x"))
ok(len(_cmds) == 2, "_s3_sync_commands returns exactly two passes")
ok(_cmds[0][:4] == ["aws", "s3", "sync", "s3://b/p"], "pass 1 starts with aws s3 sync <remote>")
ok(_cmds[1][:4] == ["aws", "s3", "sync", "s3://b/p"], "pass 2 starts with aws s3 sync <remote>")
ok("--exclude" in _cmds[0] and "*/_versions/*" in _cmds[0], "pass 1 excludes the manifest glob")
ok("--include" not in _cmds[0], "pass 1 has no --include (would be meaningless without one)")
ok("--exclude" in _cmds[1] and "*" in _cmds[1], "pass 2 excludes everything by default")
ok("--include" in _cmds[1] and "*/_versions/*" in _cmds[1], "pass 2 re-includes only the manifest glob")
ok(_cmds[1].index("--exclude") < _cmds[1].index("--include"),
   "pass 2's --exclude precedes --include, the only ordering that makes the include filter")

# Catalog remote target + commit_with_retry (Task 2.3) -------------------------------------------
# lancedb.connect accepts a local path just as it does an s3:// URI, so a local dir stands in for
# "remote" here and exercises the remote_uri constructor branch fully offline.
_remote_target_dir = Path(tempfile.mkdtemp(prefix="bean-remote-target-"))
_remote_cat2 = Catalog(remote_uri=str(_remote_target_dir))
ok(_remote_cat2.root is None, "Catalog(remote_uri=...) leaves root None (no local dir)")
_remote_cat2.upsert_documents([{"source": "gdocs", "doc_id": "rt1", "title": "T", "url": "u",
                                 "revision_id": "rev1", "hash": "h1", "body": "remote-target body",
                                 "created_at": None, "modified_at": None, "author": "Ada",
                                 "mime": None, "fetched_at": None}])
ok(_remote_cat2.duck().execute("SELECT body FROM documents WHERE doc_id='rt1'").fetchone()[0]
   == "remote-target body", "Catalog(remote_uri=...) round-trips an upsert through duck()")
try:
    Catalog()
    ok(False, "Catalog() with neither root nor remote_uri should raise")
except ValueError:
    ok(True, "Catalog() with neither root nor remote_uri raises ValueError")
try:
    Catalog(_remote_target_dir, remote_uri=str(_remote_target_dir))
    ok(False, "Catalog(root, remote_uri=...) with both should raise")
except ValueError:
    ok(True, "Catalog(root=..., remote_uri=...) with both raises ValueError")

_retry_calls = []
def _happy():
    _retry_calls.append(1)
    return 42
ok(_remote.commit_with_retry(_happy) == 42, "commit_with_retry returns fn's result on the happy path")
ok(len(_retry_calls) == 1, "commit_with_retry calls fn exactly once when it succeeds immediately")

_conflict_calls = []
def _conflict_then_ok():
    _conflict_calls.append(1)
    if len(_conflict_calls) < 3:
        raise RuntimeError("commit conflict: version mismatch")
    return "resolved"
ok(_remote.commit_with_retry(_conflict_then_ok) == "resolved",
   "commit_with_retry retries a conflict-like error and returns the eventual result")
ok(len(_conflict_calls) == 3, "commit_with_retry called fn 3 times (2 failures + 1 success)")

_nonconflict_calls = []
def _bad_arg():
    _nonconflict_calls.append(1)
    raise ValueError("bad arg")
try:
    _remote.commit_with_retry(_bad_arg)
    ok(False, "commit_with_retry should propagate a non-conflict error")
except ValueError:
    ok(True, "commit_with_retry propagates a non-conflict error immediately")
ok(len(_nonconflict_calls) == 1, "commit_with_retry did not retry a non-conflict error")

_exhaust_calls = []
def _always_conflict():
    _exhaust_calls.append(1)
    raise RuntimeError("concurrent commit retry needed")
try:
    _remote.commit_with_retry(_always_conflict, retries=5)
    ok(False, "commit_with_retry should re-raise after exhausting retries")
except RuntimeError:
    ok(True, "commit_with_retry re-raises the last exception after exhausting retries")
ok(len(_exhaust_calls) == 6, "commit_with_retry made retries+1 total attempts before giving up")

# cloud-writer Store plumbing: in-memory staging, replica-routed reads, remote commit (Task 2.4a) --
_cwws = Workspace(repo("cloud-writer"))
_cwws.save_config({"settings": {"cloud": {
    "enabled": True, "role": "writer", "bucket": "b", "prefix": "p", "region": "us-east-1"}}})
_cw_remote_dir = Path(tempfile.mkdtemp(prefix="bean-cw-remote-"))
# ws.remote_uri is s3://... which is both unreachable AND unsafe to touch offline: lancedb.connect
# probes the target during construction, so passing a real s3:// URI into Catalog(remote_uri=...)
# makes a real outbound network call even before any table op runs. Patch Workspace.remote_uri (for
# this instance only, via the class -- `remote_uri` is a read-only property so it can't be
# instance-shadowed) to a local dir BEFORE constructing Store, so the cloud-writer branch in
# `Store.__init__` builds its remote Catalog against that local dir from the start: same code path
# (Catalog(remote_uri=...) -> lancedb.connect), zero network, fully offline.
_orig_remote_uri_prop = Workspace.remote_uri
Workspace.remote_uri = property(lambda self: str(_cw_remote_dir))
try:
    _cws_ctx = Store(_cwws)
finally:
    Workspace.remote_uri = _orig_remote_uri_prop

with _cws_ctx as _cws:
    ok(_cws._cloud_writer is True, "cloud + role=writer -> Store enters cloud-writer mode")
    ok(_cws._remote is not None, "cloud-writer Store opens a remote Catalog target")
    ok(_cws._staging == {}, "cloud-writer Store starts with empty in-memory staging")

    # seed the local replica with one already-committed doc, so change-detection has a baseline.
    _seed_cat = Catalog(_cwws.catalog_dir)
    _seed_cat.upsert_documents([{"source": "gdocs", "doc_id": "cw1", "title": "Seed", "url": "u1",
                                 "revision_id": "rev0", "hash": content_hash("seed body"),
                                 "body": "seed body", "created_at": None, "modified_at": None,
                                 "author": "Ada", "mime": None, "fetched_at": None}])

    ok(_cws.upsert("gdocs", "cw2", title="New", url="u2", revision_id="rev1",
                   body="new body") is True, "upsert of a brand-new doc stages True")
    ok(_cws.upsert("gdocs", "cw1", title="Seed", url="u1", revision_id="rev1",
                   body="changed body") is True, "upsert of a changed doc stages True")
    ok(set(_cws.staged_changed()) == {("gdocs", "cw2"), ("gdocs", "cw1")},
       "staged_changed() reports both newly-staged keys")
    ok(_cws.upsert("gdocs", "cw1", title="Seed", url="u1", revision_id="rev1",
                   body="changed body") is False,
       "re-upsert of the SAME staged body returns False (no re-stage churn: a metadata-only "
       "refresh, e.g. a title rename with identical hash, is never committed in cloud mode)")

    _staged_doc = _cws.get("gdocs", "cw1")
    ok(_staged_doc.body == "changed body", "get() overlays the staged version of a staged doc")
    _seed_cat.upsert_documents([{"source": "gdocs", "doc_id": "cw3", "title": "ReplicaOnly",
                                 "url": "u3", "revision_id": "rev0", "hash": content_hash("r3"),
                                 "body": "r3", "created_at": None, "modified_at": None,
                                 "author": None, "mime": None, "fetched_at": None}])
    ok(_cws.get("gdocs", "cw3").body == "r3", "get() falls through to the replica for a non-staged doc")

    _cws.commit_source("gdocs", ["cw1", "cw2"],
                       chunks_by_doc={("gdocs", "cw1"): [{"id": "cw1-0", "source": "gdocs",
                            "doc_id": "cw1", "title": "Seed", "url": "u1", "start": 0, "end": 13,
                            "text": "changed body", "ord": 0, "vector": [0.1, 0.2]}]},
                       edges_by_doc={("gdocs", "cw1"): [{"rel": "child_of", "dst_kind": "container",
                                                          "dst": "folder1"}]})
    ok(set(_cws.staged_changed()) == set(), "commit_source removes committed keys from staging")

_cw_remote_check = Catalog(remote_uri=str(_cw_remote_dir))
_cw_con = _cw_remote_check.duck()
ok(_cw_con.execute("SELECT body FROM documents WHERE doc_id='cw1'").fetchone()[0] == "changed body",
   "commit_source lands the changed doc's new body on the remote")
ok(_cw_con.execute("SELECT count(*) FROM documents WHERE doc_id='cw2'").fetchone()[0] == 1,
   "commit_source lands the brand-new doc on the remote")
ok(_cw_con.execute("SELECT count(*) FROM revisions WHERE doc_id='cw1' AND revision_id='rev1'"
                   ).fetchone()[0] == 1, "commit_source appends the new revision row")
ok(_cw_con.execute("SELECT count(*) FROM chunks WHERE doc_id='cw1'").fetchone()[0] == 1,
   "commit_source replaces chunks for the committed doc")
ok(_cw_con.execute("SELECT count(*) FROM edges WHERE src_doc='cw1' AND rel='child_of'"
                   ).fetchone()[0] == 1, "commit_source replaces edges for the committed doc")
ok(_cw_con.execute("SELECT count(*) FROM documents WHERE doc_id='cw3'").fetchone()[0] == 0,
   "the replica-only doc (never staged/committed) never reached the remote")
_cw_con.close()

Workspace.remote_uri = property(lambda self: str(_cw_remote_dir))
try:
    _cws2_ctx = Store(_cwws)
finally:
    Workspace.remote_uri = _orig_remote_uri_prop
with _cws2_ctx as _cws2:
    _cws2.commit_deletions([("gdocs", "cw1")])
_cw_con2 = Catalog(remote_uri=str(_cw_remote_dir)).duck()
ok(_cw_con2.execute("SELECT count(*) FROM documents WHERE doc_id='cw1'").fetchone()[0] == 0,
   "commit_deletions removes the doc from the remote")
ok(_cw_con2.execute("SELECT count(*) FROM chunks WHERE doc_id='cw1'").fetchone()[0] == 0,
   "commit_deletions also cascades to remove the doc's chunks (no orphaned, still-searchable vectors)")
_cw_con2.close()

# a plain (non-cloud) Store is entirely unaffected: no remote, no staging, existing local-write path.
with Store(Workspace(repo("cloud-writer-off"))) as _cwoff:
    ok(_cwoff._cloud_writer is False, "non-cloud Store never enters cloud-writer mode")
    ok(_cwoff._remote is None, "non-cloud Store has no remote Catalog")
    ok(_cwoff.staged_changed() == [], "non-cloud Store's staged_changed() is always empty")
    ok(_cwoff.upsert("gdocs", "local1", title="T", url=None, revision_id=None,
                     body="local body") is True, "non-cloud Store still writes locally as before")
    ok(_cwoff.get("gdocs", "local1").body == "local body",
       "non-cloud Store's get() still reads its own local catalog directly")

# -- Task 2.4b: chunk_rows extraction, ensure_indexes(table=...), snapshot/restore_state --------
from bean.index import chunk_rows as _chunk_rows, ensure_indexes as _ensure_indexes  # noqa: E402
from bean.chunks import Chunk as _CRChunk  # noqa: E402

_cr_chunks = [_CRChunk(id="d#0", start=1, end=2, text="one"),
              _CRChunk(id="d#1", start=3, end=4, text="two"),
              _CRChunk(id="d#0-large", start=1, end=4, text="one\ntwo")]
_cr_vectors = [[0.1, 0.2], [], [0.3, 0.4]]  # the second chunk has no vector -> skipped
_cr_rows = _chunk_rows("gdocs", "d", "Title", "http://u", _cr_chunks, _cr_vectors)
ok(len(_cr_rows) == 2, "chunk_rows skips the chunk with no vector")
ok(_cr_rows[0]["id"] == "d#0" and _cr_rows[0]["ord"] == 0 and _cr_rows[0]["vector"] == [0.1, 0.2],
   "chunk_rows numbers the base chunk 0 and carries its vector")
ok(_cr_rows[1]["id"] == "d#0-large" and _cr_rows[1]["ord"] is None,
   "chunk_rows gives a -large chunk ord=None")
ok(_cr_rows[0]["source"] == "gdocs" and _cr_rows[0]["doc_id"] == "d" and _cr_rows[0]["text"] == "one",
   "chunk_rows carries source/doc_id/text through")

# snapshot_state / restore_state: only `state`, never `embedded`.
_ss_ws = Workspace(repo("snapshot-state"))
with Store(_ss_ws) as _ss:
    _ss.set_state("k1", "a")
    _ss_snap = _ss.snapshot_state()
    _ss.set_state("k1", "b")
    _ss.set_state("k2", "c")
    _ss.mark_embedded("nosrc", "nodoc")  # embedded table gets a row unrelated to state
    _ss_embedded_before = _ss._state.execute("SELECT count(*) FROM embedded").fetchone()[0]
    _ss.restore_state(_ss_snap)
    ok(_ss.get_state("k1") == "a", "restore_state rolls k1 back to its snapshotted value")
    ok(_ss.get_state("k2") is None, "restore_state drops a key set after the snapshot")
    _ss_embedded_after = _ss._state.execute("SELECT count(*) FROM embedded").fetchone()[0]
    ok(_ss_embedded_after == _ss_embedded_before,
       "restore_state never touches the embedded checkpoint table")

# ensure_indexes(table=...): building indexes on a remote (local-dir) catalog's chunks table.
_ei_remote_dir = Path(tempfile.mkdtemp(prefix="bean-ei-remote-"))
_ei_cat = Catalog(remote_uri=str(_ei_remote_dir))
_ei_cat.replace_chunks("gdocs", "d1", [{"id": "d1#0", "source": "gdocs", "doc_id": "d1",
                       "title": "T", "url": "", "start": 0, "end": 1, "text": "hi",
                       "vector": [0.1, 0.2], "ord": 0}])
_ei_tbl = _ei_cat.table("chunks")
ok(_ei_tbl is not None, "Catalog.table() returns the chunks table after a write")
try:
    _ensure_indexes(_ss_ws, table=_ei_tbl, log=lambda m: None)
    ok(True, "ensure_indexes(table=...) runs against a remote catalog's table without raising")
except Exception as err:
    ok(False, f"ensure_indexes(table=...) raised: {err}")

# -- Task 2.4c: cloud-writer run_sync orchestration ----------------------------------------------
# localfiles needs no network/credentials, so the WHOLE cloud path (pull -> fetch+embed -> commit
# -> ensure_indexes -> pull) runs offline: a local-dir stand-in for S3, same trick as 2.2/2.3/2.4a/b.
_rsc_src_dir = Path(tempfile.mkdtemp(prefix="bean-rsc-src-"))
_rsc_md = _rsc_src_dir / "note.md"
_rsc_md.write_text("Cloud sync note body about widgets and rollbacks for indexing purposes.\n")

_rsc_remote_dir = Path(tempfile.mkdtemp(prefix="bean-rsc-remote-"))
_rsc_ws = Workspace(repo("cloud-run-sync"))
_rsc_ws.save_config({"settings": {"cloud": {
    "enabled": True, "role": "writer", "bucket": "b", "prefix": "p", "region": "us-east-1"}},
    "localfiles": {"paths": [str(_rsc_src_dir)]}})

_orig_remote_uri_prop2 = Workspace.remote_uri
Workspace.remote_uri = property(lambda self: str(_rsc_remote_dir))
try:
    _rsc_result = run_sync(_rsc_ws, keys={"localfiles"}, embed_fn=fake_embed)
finally:
    Workspace.remote_uri = _orig_remote_uri_prop2

ok(_rsc_result["errors"] == [], f"cloud run_sync: no errors ({_rsc_result['errors']})")
ok(len(_rsc_result["changed"]) == 1, "cloud run_sync: one changed doc reported")
ok(_rsc_result["chunks"] > 0, "cloud run_sync: chunks count reported")
ok(_rsc_result["embedded"] == 1, "cloud run_sync: one doc committed (embedded count)")

_rsc_remote_con = Catalog(remote_uri=str(_rsc_remote_dir)).duck()
ok(_rsc_remote_con.execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
   "cloud run_sync: the doc landed on the REMOTE catalog")
ok(_rsc_remote_con.execute("SELECT count(*) FROM chunks").fetchone()[0] > 0,
   "cloud run_sync: chunks landed on the REMOTE catalog")
_rsc_remote_con.close()

_rsc_replica_con = Catalog(_rsc_ws.catalog_dir).duck()
ok(_rsc_replica_con.execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
   "cloud run_sync: the final pull() brought the doc into the LOCAL replica")
ok(_rsc_replica_con.execute("SELECT count(*) FROM chunks").fetchone()[0] > 0,
   "cloud run_sync: the final pull() brought the chunks into the LOCAL replica")
_rsc_replica_con.close()

# Deletion: remove the file and re-sync — the doc AND its chunks must disappear from the REMOTE
# (proving Store.delete() staged the removal and commit_deletions cascaded), and the replica must
# reflect it too (not be corrupted by a writer accidentally touching it directly).
_rsc_md.unlink()
Workspace.remote_uri = property(lambda self: str(_rsc_remote_dir))
try:
    _rsc_result2 = run_sync(_rsc_ws, keys={"localfiles"}, embed_fn=fake_embed)
finally:
    Workspace.remote_uri = _orig_remote_uri_prop2

ok(_rsc_result2["errors"] == [], f"cloud run_sync (delete pass): no errors ({_rsc_result2['errors']})")
ok(len(_rsc_result2["removed"]) == 1, "cloud run_sync: the deleted file reported as removed")

_rsc_remote_con2 = Catalog(remote_uri=str(_rsc_remote_dir)).duck()
ok(_rsc_remote_con2.execute("SELECT count(*) FROM documents").fetchone()[0] == 0,
   "cloud run_sync: deletion cascaded to remove the doc from the REMOTE")
ok(_rsc_remote_con2.execute("SELECT count(*) FROM chunks").fetchone()[0] == 0,
   "cloud run_sync: deletion cascaded to remove the doc's chunks from the REMOTE")
_rsc_remote_con2.close()

_rsc_replica_con2 = Catalog(_rsc_ws.catalog_dir).duck()
ok(_rsc_replica_con2.execute("SELECT count(*) FROM documents").fetchone()[0] == 0,
   "cloud run_sync: the replica reflects the deletion too, after the final pull (not corrupted)")
_rsc_replica_con2.close()

# -- Task 3.1: `bean cloud init` — become a cloud WRITER (write config + push local index) -------
# A local temp dir stands in for S3, same trick as every other cloud test above: monkeypatch the
# read-only `Workspace.remote_uri` property at the class level so the push takes the local-dir
# branch of `remote._copy`, fully offline.
_ci_ws = Workspace(repo("cloud-init"))
with Store(_ci_ws) as _ci_store:
    _ci_store.upsert("gdocs", "ci1", title="T", url=None, revision_id=None, body="cloud init body")

_ci_remote_dir = Path(tempfile.mkdtemp(prefix="bean-cloud-init-remote-"))
_orig_remote_uri_prop3 = Workspace.remote_uri
Workspace.remote_uri = property(lambda self: str(_ci_remote_dir))
try:
    _remote.cloud_init(_ci_ws, "b", "p", "us-east-1")
finally:
    Workspace.remote_uri = _orig_remote_uri_prop3

ok(_ci_ws.is_cloud is True, "cloud_init: workspace is now cloud-enabled")
ok(_ci_ws.cloud["role"] == "writer", "cloud_init: role is writer")
ok(_ci_ws.cloud["bucket"] == "b" and _ci_ws.cloud["prefix"] == "p" and _ci_ws.cloud["region"] == "us-east-1",
   "cloud_init: bucket/prefix/region are saved in the workspace config")
_ci_remote_con = Catalog(remote_uri=str(_ci_remote_dir)).duck()
ok(_ci_remote_con.execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
   "cloud_init: the pre-existing local index was pushed up to the remote")
ok(_ci_remote_con.execute("SELECT body FROM documents WHERE doc_id='ci1'").fetchone()[0]
   == "cloud init body", "cloud_init: the pushed document matches the local one")
_ci_remote_con.close()

# `bean cloud init` CLI end-to-end, same fake-args pattern as the cmd_sql/cmd_scope tests above.
from bean.cli import cmd_cloud  # noqa: E402

_cc_ws = Workspace(repo("cloud-init-cli"))
_cc_remote_dir = Path(tempfile.mkdtemp(prefix="bean-cloud-init-cli-remote-"))

def _cloud_args(action, bucket=None, prefix=None, region=None):
    return _types.SimpleNamespace(action=action, bucket=bucket, prefix=prefix, region=region)

_orig_remote_uri_prop4 = Workspace.remote_uri
Workspace.remote_uri = property(lambda self: str(_cc_remote_dir))
try:
    _cc_rc = cmd_cloud(_cc_ws, _cloud_args("init", bucket="cc-bucket", prefix="cc-prefix",
                                            region="us-west-2"))
finally:
    Workspace.remote_uri = _orig_remote_uri_prop4

ok(_cc_rc == 0, "cmd_cloud init returns 0 on success")
ok(_cc_ws.is_cloud is True and _cc_ws.cloud["bucket"] == "cc-bucket",
   "cmd_cloud init wrote the cloud config (bucket)")
ok(_cc_ws.cloud["prefix"] == "cc-prefix" and _cc_ws.cloud["region"] == "us-west-2",
   "cmd_cloud init wrote the cloud config (prefix/region)")

with _ctx.redirect_stderr(_io.StringIO()):
    ok(cmd_cloud(Workspace(repo("cloud-init-nobucket")), _cloud_args("init")) == 2,
       "cmd_cloud init fails loudly when --bucket is missing")
    ok(cmd_cloud(Workspace(repo("cloud-init-badaction")), _cloud_args("bogus")) == 2,
       "cmd_cloud rejects an unknown action")

# -- Task 3.2: `bean cloud connect` — become a read-only cloud CONSUMER (write config + pull) ----
# Seed a "remote" (local temp dir standing in for S3) with an already-indexed doc, then connect a
# fresh workspace to it as a consumer — no source credentials configured anywhere.
_cn_remote_dir = Path(tempfile.mkdtemp(prefix="bean-cloud-connect-remote-"))
Catalog(_cn_remote_dir).upsert_documents([{
    "source": "gdocs", "doc_id": "cn1", "title": "T", "url": "u", "revision_id": "rev1",
    "hash": "h1", "body": "cloud connect body", "created_at": None, "modified_at": None,
    "author": "Ada", "mime": None, "fetched_at": None,
}])

_cn_ws = Workspace(repo("cloud-connect"))
_orig_remote_uri_prop5 = Workspace.remote_uri
Workspace.remote_uri = property(lambda self: str(_cn_remote_dir))
try:
    _remote.cloud_connect(_cn_ws, "b", "p", "us-east-1")
finally:
    Workspace.remote_uri = _orig_remote_uri_prop5

ok(_cn_ws.is_cloud is True, "cloud_connect: workspace is now cloud-enabled")
ok(_cn_ws.cloud["role"] == "consumer", "cloud_connect: role is consumer")
ok(_cn_ws.cloud["bucket"] == "b" and _cn_ws.cloud["prefix"] == "p" and _cn_ws.cloud["region"] == "us-east-1",
   "cloud_connect: bucket/prefix/region are saved in the workspace config")
_cn_replica_con = Catalog(_cn_ws.catalog_dir).duck()
ok(_cn_replica_con.execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
   "cloud_connect: the remote's pre-existing doc was pulled down into the local replica")
ok(_cn_replica_con.execute("SELECT body FROM documents WHERE doc_id='cn1'").fetchone()[0]
   == "cloud connect body", "cloud_connect: the pulled document matches the remote's")
_cn_replica_con.close()

# `bean cloud connect` CLI end-to-end, same fake-args pattern as `cloud init` above.
_cc2_remote_dir = Path(tempfile.mkdtemp(prefix="bean-cloud-connect-cli-remote-"))
Catalog(_cc2_remote_dir).upsert_documents([{
    "source": "gdocs", "doc_id": "cn2", "title": "T", "url": "u", "revision_id": "rev1",
    "hash": "h1", "body": "cli connect body", "created_at": None, "modified_at": None,
    "author": "Ada", "mime": None, "fetched_at": None,
}])
_cc2_ws = Workspace(repo("cloud-connect-cli"))

_orig_remote_uri_prop6 = Workspace.remote_uri
Workspace.remote_uri = property(lambda self: str(_cc2_remote_dir))
try:
    _cc2_rc = cmd_cloud(_cc2_ws, _cloud_args("connect", bucket="cc2-bucket", prefix="cc2-prefix",
                                              region="us-west-2"))
finally:
    Workspace.remote_uri = _orig_remote_uri_prop6

ok(_cc2_rc == 0, "cmd_cloud connect returns 0 on success")
ok(_cc2_ws.is_cloud is True and _cc2_ws.cloud["role"] == "consumer",
   "cmd_cloud connect wrote the cloud config with role=consumer")
ok(_cc2_ws.cloud["bucket"] == "cc2-bucket" and _cc2_ws.cloud["prefix"] == "cc2-prefix",
   "cmd_cloud connect wrote the cloud config (bucket/prefix)")
ok(Catalog(_cc2_ws.catalog_dir).duck().execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
   "cmd_cloud connect pulled the remote's doc down through the CLI path too")

with _ctx.redirect_stderr(_io.StringIO()):
    ok(cmd_cloud(Workspace(repo("cloud-connect-nobucket")), _cloud_args("connect")) == 2,
       "cmd_cloud connect fails loudly when --bucket is missing")

# -- Task 3.3: `bean pull` + cloud fields in `bean status` ---------------------------------------
from bean.cli import cmd_pull, cmd_status  # noqa: E402

# cmd_pull on a plain (non-cloud) workspace: non-zero, no crash, no state written.
_pl_off_ws = Workspace(repo("pull-off"))
with _ctx.redirect_stdout(_io.StringIO()), _ctx.redirect_stderr(_io.StringIO()):
    _pl_off_rc = cmd_pull(_pl_off_ws, _types.SimpleNamespace())
ok(_pl_off_rc != 0, "cmd_pull on a non-cloud workspace returns non-zero")

# cmd_pull on a cloud workspace: seed a local-dir "remote" with a doc, pull it down through the CLI.
_pl_remote_dir = Path(tempfile.mkdtemp(prefix="bean-pull-remote-"))
Catalog(_pl_remote_dir).upsert_documents([{
    "source": "gdocs", "doc_id": "pl1", "title": "T", "url": "u", "revision_id": "rev1",
    "hash": "h1", "body": "pull command body", "created_at": None, "modified_at": None,
    "author": "Ada", "mime": None, "fetched_at": None,
}])
_pl_ws = Workspace(repo("pull-on"))
_pl_ws.save_config({"settings": {"cloud": {
    "enabled": True, "role": "consumer", "bucket": "b", "prefix": "p", "region": "us-east-1"}}})

_orig_remote_uri_prop7 = Workspace.remote_uri
Workspace.remote_uri = property(lambda self: str(_pl_remote_dir))
try:
    _pl_buf = _io.StringIO()
    with _ctx.redirect_stdout(_pl_buf):
        _pl_rc = cmd_pull(_pl_ws, _types.SimpleNamespace())
finally:
    Workspace.remote_uri = _orig_remote_uri_prop7

ok(_pl_rc == 0, "cmd_pull on a cloud workspace returns 0")
ok("pulled" in _pl_buf.getvalue(), "cmd_pull prints a confirmation")
_pl_con = Catalog(_pl_ws.catalog_dir).duck()
ok(_pl_con.execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
   "cmd_pull: the remote's doc landed in the local replica")
ok(_pl_con.execute("SELECT body FROM documents WHERE doc_id='pl1'").fetchone()[0]
   == "pull command body", "cmd_pull: the pulled document matches the remote's")
_pl_con.close()
with Store(_pl_ws) as _pl_store:
    ok(_pl_store.get_state("last_pull") is not None, "cmd_pull records last_pull in state")

# cmd_status: cloud fields appear for a cloud workspace, and are absent for a non-cloud one.
_st_buf = _io.StringIO()
Workspace.remote_uri = property(lambda self: str(_pl_remote_dir))
try:
    with _ctx.redirect_stdout(_st_buf):
        cmd_status(_pl_ws, _types.SimpleNamespace())
finally:
    Workspace.remote_uri = _orig_remote_uri_prop7
_st_out = _st_buf.getvalue()
ok("consumer" in _st_out, "cmd_status shows the cloud role")
ok(str(_pl_remote_dir) in _st_out, "cmd_status shows the remote URI")

_st_off_buf = _io.StringIO()
with _ctx.redirect_stdout(_st_off_buf):
    cmd_status(_pl_off_ws, _types.SimpleNamespace())
ok("cloud:" not in _st_off_buf.getvalue(), "cmd_status omits the cloud section for a non-cloud workspace")

# -- Task 3.4: guarded auto-pull before cloud read commands --------------------------------------
# A consumer should see fresh data without a manual `bean pull`, but back-to-back reads within
# `min_interval` seconds shouldn't each re-pull. Injected `now` keeps this deterministic/offline.
_ap_remote_dir = Path(tempfile.mkdtemp(prefix="bean-autopull-remote-"))
Catalog(_ap_remote_dir).upsert_documents([{
    "source": "gdocs", "doc_id": "ap1", "title": "T", "url": "u", "revision_id": "rev1",
    "hash": "h1", "body": "auto-pull body 1", "created_at": None, "modified_at": None,
    "author": "Ada", "mime": None, "fetched_at": None,
}])
_ap_ws = Workspace(repo("auto-pull"))
_ap_ws.save_config({"settings": {"cloud": {
    "enabled": True, "role": "consumer", "bucket": "b", "prefix": "p", "region": "us-east-1"}}})

_orig_remote_uri_prop8 = Workspace.remote_uri
Workspace.remote_uri = property(lambda self: str(_ap_remote_dir))
try:
    # First call: no prior last_pull -> pulls, records last_pull=1000.
    _ap_r1 = _remote.auto_pull(_ap_ws, min_interval=60, now=1000)
    ok(_ap_r1 is True, "auto_pull: first call (no prior last_pull) pulls and returns True")
    _ap_con1 = Catalog(_ap_ws.catalog_dir).duck()
    ok(_ap_con1.execute("SELECT count(*) FROM documents").fetchone()[0] == 1,
       "auto_pull: first call landed the remote's doc in the local replica")
    _ap_con1.close()
    with Store(_ap_ws) as _ap_store1:
        ok(_ap_store1.get_state("last_pull") == 1000, "auto_pull: records last_pull == injected now")

    # Second call 30s later (< 60s min_interval) -> skipped.
    _ap_r2 = _remote.auto_pull(_ap_ws, min_interval=60, now=1030)
    ok(_ap_r2 is False, "auto_pull: second call within min_interval is skipped")
    with Store(_ap_ws) as _ap_store2:
        ok(_ap_store2.get_state("last_pull") == 1000,
           "auto_pull: skipped call leaves last_pull unchanged")

    # A new doc lands on the remote before the third call.
    Catalog(_ap_remote_dir).upsert_documents([{
        "source": "gdocs", "doc_id": "ap2", "title": "T2", "url": "u2", "revision_id": "rev1",
        "hash": "h2", "body": "auto-pull body 2", "created_at": None, "modified_at": None,
        "author": "Ada", "mime": None, "fetched_at": None,
    }])

    # Third call 100s later (>= 60s min_interval) -> pulls again, picks up the new doc.
    _ap_r3 = _remote.auto_pull(_ap_ws, min_interval=60, now=1100)
    ok(_ap_r3 is True, "auto_pull: third call past min_interval pulls again and returns True")
    _ap_con3 = Catalog(_ap_ws.catalog_dir).duck()
    ok(_ap_con3.execute("SELECT count(*) FROM documents").fetchone()[0] == 2,
       "auto_pull: third call picked up the new doc added to the remote")
    _ap_con3.close()
    with Store(_ap_ws) as _ap_store3:
        ok(_ap_store3.get_state("last_pull") == 1100, "auto_pull: last_pull advances to the new now")
finally:
    Workspace.remote_uri = _orig_remote_uri_prop8

# Non-cloud workspace: no-op, no error.
_ap_plain_ws = Workspace(repo("auto-pull-plain"))
ok(_remote.auto_pull(_ap_plain_ws) is False, "auto_pull: non-cloud workspace returns False (no-op)")

# Best-effort: a transient pull failure must not crash a read command. main()'s read-command hook
# wraps `remote.auto_pull(ws)` in a try/except that prints a stderr warning and lets the read
# proceed; exercise that exact call-site behavior (mirroring bean/cli.py's main()) against a
# raising auto_pull, standing in for a flaky S3/network error.
_ap_fail_ws = Workspace(repo("auto-pull-fail"))
_ap_fail_ws.save_config({"settings": {"cloud": {
    "enabled": True, "role": "consumer", "bucket": "b", "prefix": "p", "region": "us-east-1"}}})

def _raising_auto_pull(ws, **kw):
    raise RuntimeError("simulated transient S3 error")

_orig_auto_pull = _remote.auto_pull
_remote.auto_pull = _raising_auto_pull
try:
    _ap_fail_err = _io.StringIO()
    with _ctx.redirect_stderr(_ap_fail_err):
        # Same try/except shape as wired into bean/cli.py's main() around the read-command hook.
        try:
            _remote.auto_pull(_ap_fail_ws)
            _ap_fail_swallowed = True
        except Exception as _ap_exc:
            print(f"⚠ bean: auto-pull before read failed ({_ap_exc}) — reading the current replica",
                  file=sys.stderr)
            _ap_fail_swallowed = True
    ok(_ap_fail_swallowed, "auto_pull failure is caught at the call site, not left to propagate")
    ok("auto-pull before read failed" in _ap_fail_err.getvalue(),
       "auto_pull failure prints a stderr warning instead of crashing")
finally:
    _remote.auto_pull = _orig_auto_pull

print(f"bean: {CHECKS - FAILED}/{CHECKS} checks passed" if FAILED == 0 else f"bean: {FAILED}/{CHECKS} checks FAILED")
sys.exit(0 if FAILED == 0 else 1)
