"""Text extraction for office "docs". Word (.docx/.docm), OpenDocument (.odt/.fodt) and RTF are
read with stdlib alone (they're just zipped XML). Presentations (.pptx) and spreadsheets (.xlsx)
go through purpose-built libraries — python-pptx and openpyxl — because reassembling slide/cell
structure by hand is fragile; both are lazy-imported so they only matter when such a file appears.
Each extractor returns plain text with structure (paragraphs, slides, sheets) preserved, or raises
so the caller can log-skip one unreadable file without aborting the whole sync.

Legacy binary formats (.doc, .ppt) are intentionally not handled — they need an external converter
(antiword/LibreOffice); the caller skips them with a clear message."""

from __future__ import annotations

import zipfile
from pathlib import Path

DOCX_EXT = {".docx", ".docm"}
ODT_EXT = {".odt", ".fodt"}
RTF_EXT = {".rtf"}
PPTX_EXT = {".pptx"}
XLSX_EXT = {".xlsx"}
OFFICE_EXT = DOCX_EXT | ODT_EXT | RTF_EXT | PPTX_EXT | XLSX_EXT

# OOXML / ODF namespaces we pull text out of.
_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
_ODT_TEXT_TAGS = {"p", "h"}  # text:p (paragraph) and text:h (heading), namespace-stripped


def extract_office(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in DOCX_EXT:
        return _docx(path)
    if ext in ODT_EXT:
        return _odt(path)
    if ext in RTF_EXT:
        return _rtf(path.read_bytes().decode("latin-1", errors="replace"))
    if ext in PPTX_EXT:
        return _pptx(path)
    if ext in XLSX_EXT:
        return _xlsx(path)
    raise ValueError(f"unsupported office extension: {ext}")


def _pptx(path: Path) -> str:
    """PowerPoint .pptx: text from every shape (and table) per slide, plus speaker notes, kept in
    slide order with a heading per slide so a chunk still says which slide it came from."""
    try:
        from pptx import Presentation
    except ImportError as err:
        raise RuntimeError("reading .pptx needs `pip install python-pptx`") from err
    out = []
    for i, slide in enumerate(Presentation(str(path)).slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append("\t".join(c.text for c in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append(f"(notes) {slide.notes_slide.notes_text_frame.text}")
        body = "\n".join(p for p in parts if p and p.strip())
        out.append(f"## Slide {i}\n{body}" if body else f"## Slide {i}")
    return "\n\n".join(out).strip()


def _xlsx(path: Path) -> str:
    """Excel .xlsx: each sheet's non-empty rows as tab-joined cells, computed values (not formulas),
    under a heading per sheet. read_only streaming so a large workbook doesn't load fully into RAM."""
    try:
        from openpyxl import load_workbook
    except ImportError as err:
        raise RuntimeError("reading .xlsx needs `pip install openpyxl`") from err
    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        out = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                out.append(f"## {sheet.title}\n" + "\n".join(rows))
        return "\n\n".join(out).strip()
    finally:
        wb.close()


def _docx(path: Path) -> str:
    """Word .docx/.docm: text lives in word/document.xml as runs (<w:t>) grouped by paragraph
    (<w:p>). We join runs within a paragraph and paragraphs with newlines; tabs and breaks become
    whitespace."""
    import xml.etree.ElementTree as ET

    with zipfile.ZipFile(path) as z:
        root = ET.fromstring(z.read("word/document.xml"))
    paras = []
    for p in root.iter(f"{_W}p"):
        buf = []
        for node in p.iter():
            tag = node.tag
            if tag == f"{_W}t" and node.text:
                buf.append(node.text)
            elif tag == f"{_W}tab":
                buf.append("\t")
            elif tag in (f"{_W}br", f"{_W}cr"):
                buf.append("\n")
        paras.append("".join(buf))
    return "\n".join(paras).strip()


def _odt(path: Path) -> str:
    """OpenDocument .odt: text lives in content.xml. We take each paragraph/heading element's full
    inner text (spans, links, etc. flatten via itertext)."""
    import xml.etree.ElementTree as ET

    with zipfile.ZipFile(path) as z:
        root = ET.fromstring(z.read("content.xml"))
    paras = []
    for node in root.iter():
        if node.tag.split("}")[-1] in _ODT_TEXT_TAGS:
            paras.append("".join(node.itertext()))
    return "\n".join(paras).strip()


def _rtf(text: str) -> str:
    """Crude RTF → text: drop the font/color/stylesheet groups, decode \\'hh and unicode \\uN
    escapes, turn \\par into newlines, then strip remaining control words and braces. Not a full
    parser — good enough to make RTF searchable."""
    import re

    # Drop groups we never want as body text (font tables, color tables, stylesheets, info).
    text = re.sub(r"\{\\\*?\\(?:fonttbl|colortbl|stylesheet|info|generator)[^{}]*\}", " ", text)
    # Paragraph/line/tab breaks: a control word swallows one trailing delimiter space.
    text = re.sub(r"\\par\b ?", "\n", text)
    text = re.sub(r"\\line\b ?", "\n", text)
    text = re.sub(r"\\tab\b ?", "\t", text)
    text = re.sub(r"\\'([0-9a-fA-F]{2})", lambda m: bytes.fromhex(m.group(1)).decode("latin-1"), text)
    text = re.sub(r"\\u(-?\d+)\??", lambda m: chr(int(m.group(1)) % 0x10000), text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)  # remaining control words
    text = text.replace("{", "").replace("}", "").replace("\\", "")
    return re.sub(r"[ \t]+\n", "\n", text).strip()
